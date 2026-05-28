from pptx import Presentation
prs = Presentation(r'c:\Users\leela\compliance-ai-saas\AWS_Compliance_AI_SaaS_AWSonly_FullPlan.pptx')
print(f'Total slides: {len(prs.slides)}')
print(f'Slide size: {prs.slide_width.inches:.2f}" x {prs.slide_height.inches:.2f}"')
print()
for i, slide in enumerate(prs.slides, 1):
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t and len(t) > 2:
                    texts.append(t)
    title = texts[0] if texts else '(no text)'
    all_text = ' | '.join(texts[:4])
    print(f'Slide {i:2d}: {all_text[:120]}'.encode('ascii', 'replace').decode())
