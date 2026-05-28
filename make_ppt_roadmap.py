# VigiliCloud - Internal Roadmap & Status PPT
# Run: C:\Users\leela\anaconda3\python.exe make_ppt_roadmap.py

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Palette
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY    = RGBColor(0xF1, 0xF5, 0xF9)
BORDER   = RGBColor(0xE2, 0xE8, 0xF0)
NAVY     = RGBColor(0x0F, 0x27, 0x6D)
NAVYMD   = RGBColor(0x1D, 0x4E, 0xD8)
TEAL     = RGBColor(0x05, 0x96, 0x69)
LTEAL    = RGBColor(0x10, 0xB9, 0x81)
TEAL_BG  = RGBColor(0xEC, 0xFD, 0xF5)
NAVY_BG  = RGBColor(0xEF, 0xF6, 0xFF)
TEXT     = RGBColor(0x0F, 0x17, 0x2A)
SUBTEXT  = RGBColor(0x47, 0x55, 0x69)
MUTED    = RGBColor(0x94, 0xA3, 0xB8)
RED      = RGBColor(0xDC, 0x26, 0x26)
RED_BG   = RGBColor(0xFE, 0xF2, 0xF2)
ORANGE   = RGBColor(0xEA, 0x58, 0x0C)
ORG_BG   = RGBColor(0xFF, 0xF7, 0xED)
AMBER    = RGBColor(0xD9, 0x77, 0x06)
AMB_BG   = RGBColor(0xFF, 0xFB, 0xEB)
GREEN    = RGBColor(0x05, 0x96, 0x69)
GREEN_BG = RGBColor(0xEC, 0xFD, 0xF5)
PURPLE   = RGBColor(0x7C, 0x3A, 0xED)
PURP_BG  = RGBColor(0xF5, 0xF3, 0xFF)
SLATE    = RGBColor(0x64, 0x74, 0x8B)
SLT_BG   = RGBColor(0xF8, 0xFA, 0xFC)

STATUS_COLORS = {
    "DONE":        (TEAL,   TEAL_BG),
    "IN PROGRESS": (AMBER,  AMB_BG),
    "NEXT":        (NAVYMD, NAVY_BG),
    "PLANNED":     (SLATE,  SLT_BG),
    "FUTURE":      (PURPLE, PURP_BG),
    "BLOCKED":     (RED,    RED_BG),
}

TOTAL = 10

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def bg(slide, color=LGRAY):
    fill = slide.background.fill
    fill.solid(); fill.fore_color.rgb = color


def rect(slide, l, t, w, h, fill_color, border_color=None, border_w=0.75):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill_color
    if border_color:
        s.line.color.rgb = border_color; s.line.width = Pt(border_w)
    else:
        s.line.fill.background()
    return s


def txt(slide, text, l, t, w, h, size=12, bold=False, color=TEXT,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return tb


def hdr(slide, title, subtitle=None, accent=TEAL):
    rect(slide, 0, 0, 13.33, 1.3, NAVY)
    rect(slide, 0, 1.3, 13.33, 0.045, accent)
    txt(slide, title, 0.45, 0.1, 10.5, 0.65, size=23, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, 0.45, 0.77, 11.8, 0.42, size=11, color=RGBColor(0xBF, 0xDB, 0xFF))


def slide_num(slide, n):
    txt(slide, f"{n} of {TOTAL}", 12.2, 7.1, 0.9, 0.3, size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def card(slide, l, t, w, h, fill=WHITE, border=BORDER, bw=0.75):
    return rect(slide, l, t, w, h, fill, border, bw)


# ===================================================================
# SLIDE 1 - COVER
# ===================================================================
s = prs.slides.add_slide(BLANK)
bg(s, NAVY)
rect(s, 0, 0, 0.08, 7.5, TEAL)
rect(s, 0, 7.42, 13.33, 0.08, TEAL)

logo = rect(s, 0.6, 0.55, 1.0, 1.0, TEAL)
logo.line.fill.background()
txt(s, "V", 0.6, 0.55, 1.0, 1.0, size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(s, "VigiliCloud", 1.8, 0.6, 7, 0.7, size=34, bold=True, color=WHITE)
txt(s, "Internal Roadmap & Build Status", 1.8, 1.35, 7, 0.38, size=14, color=LTEAL)

rect(s, 0.6, 1.88, 9.5, 0.05, TEAL)

txt(s, "What's done. What's next. What Phase 2 (MCP) means.",
    0.6, 2.1, 9.5, 0.48, size=17, bold=True, color=WHITE)
txt(s, "Honest internal status across all planned features -- May 2026",
    0.6, 2.65, 9.5, 0.35, size=12, color=RGBColor(0xBF, 0xDB, 0xFF))

legend_items = [
    ("DONE",        TEAL,   "Feature is live on production"),
    ("IN PROGRESS", AMBER,  "Partially built, needs completion"),
    ("NEXT",        NAVYMD, "Highest priority to build next"),
    ("PLANNED",     SLATE,  "Confirmed for roadmap, not started"),
    ("FUTURE",      PURPLE, "Phase 2 -- after product-market fit"),
    ("BLOCKED",     RED,    "Needs external action to unblock"),
]
txt(s, "Status Key:", 0.6, 3.3, 2.0, 0.32, size=11, bold=True, color=MUTED)
for i, (label, col, meaning) in enumerate(legend_items):
    lx = 0.6 + (i % 2) * 5.5
    ly = 3.7 + (i // 2) * 0.72
    dot = rect(s, lx, ly + 0.08, 0.28, 0.28, col)
    dot.line.fill.background()
    txt(s, label, lx + 0.38, ly, 1.5, 0.3, size=11, bold=True, color=col)
    txt(s, meaning, lx + 0.38, ly + 0.32, 4.8, 0.28, size=9.5, color=MUTED)

txt(s, "VigiliCloud  |  INTERNAL USE ONLY  |  Not for distribution",
    0.6, 6.9, 12.0, 0.35, size=10, color=MUTED, align=PP_ALIGN.CENTER, italic=True)

slide_num(s, 1)

# ===================================================================
# SLIDE 2 - OVERALL PROGRESS SNAPSHOT
# ===================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
hdr(s, "Overall Progress Snapshot", "Where VigiliCloud stands today -- May 2026")

categories = [
    ("Core Product (Scanner + AI + Exports)", 9, 9, TEAL,   "100%", "All MVP features are live on app.vigilicloud.com"),
    ("Billing (Razorpay)",                    7, 8, AMBER,  "87%",  "Code 100% done -- only Render env vars missing"),
    ("Notifications (Email Alerts)",          7, 8, TEAL,   "87%",  "Email live via Resend -- domain email pending"),
    ("Advanced Features (Tickets/Drift/API)", 0, 8, NAVYMD, "0%",   "Q3 2026 roadmap -- not MVP, expansion features"),
    ("MSP Multi-Tenant Mode",                 0, 6, SLATE,  "0%",   "Q4 2026 -- needed for MSP plan at scale"),
    ("Phase 2 MCP Governance",                0, 6, PURPLE, "0%",   "2027+ -- the powerful Phase 2 vision"),
]

for i, (cat, done, total, col, pct, note) in enumerate(categories):
    y = 1.6 + i * 0.92
    card(s, 0.45, y, 12.43, 0.82, fill=WHITE, border=BORDER)
    rect(s, 0.45, y, 0.06, 0.82, col)

    txt(s, cat, 0.65, y + 0.06, 3.5, 0.3, size=11, bold=True, color=TEXT)
    txt(s, note, 0.65, y + 0.44, 3.5, 0.28, size=9.5, color=SUBTEXT, italic=True)

    bar_x = 4.35; bar_y = y + 0.26; bar_w = 6.8; bar_h = 0.3
    rect(s, bar_x, bar_y, bar_w, bar_h, LGRAY, BORDER, 0.5)
    filled = bar_w * (done / total) if total > 0 else 0
    if filled > 0:
        filled_bar = rect(s, bar_x, bar_y, filled, bar_h, col)
        filled_bar.line.fill.background()

    txt(s, pct, 11.25, y + 0.2, 1.45, 0.42, size=20, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(s, f"{done}/{total} done", 11.25, y + 0.56, 1.45, 0.22, size=8.5, color=MUTED, align=PP_ALIGN.CENTER)

card(s, 0.45, 7.12, 12.43, 0.28, fill=TEAL_BG, border=TEAL, bw=1.0)
txt(s, "MVP is live and working. Billing needs 1 action (30 min). Advanced features are the growth roadmap -- not missing MVP requirements.",
    0.65, 7.16, 12.0, 0.2, size=10, bold=True, color=TEAL)

slide_num(s, 2)

# ===================================================================
# SLIDE 3 - WHAT'S DONE (CORE MVP)
# ===================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
hdr(s, "What's Done -- Core MVP (100% Live)", "All of these are running right now on app.vigilicloud.com", accent=TEAL)

done_items = [
    ("AWS Account Onboarding",     "Read-only IAM role + STS AssumeRole. ARN input, connection test, account listing. Zero credentials stored ever."),
    ("10 Security Checks",         "S3 public access, Root keys, IAM permissions, MFA, Security Groups, RDS, EBS, CloudTrail, VPC Flow Logs, KMS rotation."),
    ("Parallel Scan Engine",       "All 10 checks run simultaneously via boto3. Severity classification (CRITICAL/HIGH/MEDIUM). Per-resource findings in database."),
    ("Dashboard & Findings UI",    "Findings table with severity/service/status filters. FindingDetail panel with exact fix guidance for every finding."),
    ("Claude AI Executive Summary","POST /scans/{id}/ai-analysis calls Claude Haiku. Generates plain-English executive summary with prioritised fixes."),
    ("CSV / JSON Export",          "One-click compliance evidence download from scans page. All findings with resource IDs, severity, remediation steps."),
    ("Email Alerts on CRITICAL",   "Resend API sends immediate email when CRITICAL findings are found. Configured via RESEND_API_KEY + FROM_EMAIL."),
    ("Auth System (Session-based)","Signup, signin, signout, /auth/me. Session cookies with TTL. Rate limiting on all auth endpoints. Admin seed on startup."),
    ("Production Deployment",      "Backend (FastAPI on Render) + UI (Next.js on Render). Auto-deploy from main branch GitHub pushes. Postgres database."),
]
for i, (name, desc) in enumerate(done_items):
    y = 1.55 + i * 0.62
    card(s, 0.45, y, 12.43, 0.56, fill=TEAL_BG, border=TEAL, bw=0.5)
    rect(s, 0.45, y, 0.06, 0.56, TEAL)
    tick = rect(s, 0.6, y + 0.13, 0.28, 0.28, TEAL)
    tick.line.fill.background()
    txt(s, "LIVE", 0.6, y + 0.15, 0.28, 0.24, size=5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, name, 1.0, y + 0.04, 3.1, 0.28, size=10.5, bold=True, color=TEAL)
    txt(s, desc, 1.0, y + 0.3, 11.7, 0.24, size=9, color=SUBTEXT)

slide_num(s, 3)

# ===================================================================
# SLIDE 4 - RAZORPAY STATUS (THE FULL TRUTH)
# ===================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
hdr(s, "Razorpay Billing -- 87% Done (Code Complete)", "The code is 100% written. Only 6 environment variables are missing.", accent=AMBER)

txt(s, "What's already in the code (billing.py + deps.py):", 0.45, 1.52, 7.8, 0.35, size=13, bold=True, color=NAVY)

code_done = [
    ("Subscription creation",     "POST /billing/subscribe -- creates Razorpay subscription for starter/pro/MSP"),
    ("Payment verification",      "POST /billing/verify-payment -- verifies HMAC signature after Razorpay checkout flow"),
    ("Webhook handler",           "POST /billing/webhook -- handles subscription.activated, halted, payment.captured events"),
    ("Cancel subscription",       "POST /billing/cancel -- cancels at end of billing cycle, resets user plan to free"),
    ("Live status sync",          "GET /billing/sync -- fetches live subscription status directly from Razorpay API"),
    ("Signature verification",    "verify_razorpay_payment_signature() and verify_razorpay_webhook_signature() both implemented"),
    ("Plan ID configuration",     "RAZORPAY_PLAN_STARTER/PRO/MSP read from env vars -- plan IDs come from Razorpay dashboard"),
    ("Billing status endpoint",   "GET /billing/status -- returns plan tier, account limit, sub ID, Razorpay config state"),
]
for i, (name, desc) in enumerate(code_done):
    y = 1.95 + i * 0.54
    card(s, 0.45, y, 7.8, 0.48, fill=TEAL_BG, border=TEAL, bw=0.5)
    rect(s, 0.45, y, 0.05, 0.48, TEAL)
    txt(s, name, 0.62, y + 0.04, 2.5, 0.22, size=9.5, bold=True, color=TEAL)
    txt(s, desc, 0.62, y + 0.26, 7.52, 0.2, size=8.5, color=SUBTEXT)

# Right panel -- what's needed
txt(s, "The only remaining step:", 8.45, 1.52, 4.73, 0.35, size=13, bold=True, color=NAVY)
txt(s, "Add these 6 env vars to Render:", 8.45, 1.88, 4.73, 0.28, size=10.5, color=SUBTEXT)

missing_vars = [
    ("RAZORPAY_KEY_ID",        "API Key ID from Razorpay Dashboard"),
    ("RAZORPAY_KEY_SECRET",    "API Key Secret from Razorpay Dashboard"),
    ("RAZORPAY_WEBHOOK_SECRET","Webhook secret from Razorpay -> Webhooks"),
    ("RAZORPAY_PLAN_STARTER",  "Plan ID for the Starter plan"),
    ("RAZORPAY_PLAN_PRO",      "Plan ID for the Pro plan"),
    ("RAZORPAY_PLAN_MSP",      "Plan ID for the MSP plan"),
]
for i, (var, desc) in enumerate(missing_vars):
    y = 2.25 + i * 0.82
    card(s, 8.45, y, 4.73, 0.72, fill=AMB_BG, border=AMBER, bw=0.75)
    txt(s, var, 8.62, y + 0.05, 4.38, 0.3, size=9.5, bold=True, color=AMBER)
    txt(s, desc, 8.62, y + 0.38, 4.38, 0.26, size=9, color=SUBTEXT)

card(s, 8.45, 7.22, 4.73, 0.15, fill=TEAL_BG, border=TEAL, bw=0.75)

txt(s, "Steps: Razorpay Dashboard -> Create Plans (Starter/Pro/MSP) -> Copy Plan IDs -> Render Dashboard -> Backend Service -> Environment -> Add Variables",
    0.45, 7.12, 7.82, 0.3, size=9, color=SUBTEXT, italic=True)

card(s, 8.45, 6.9, 4.73, 0.25, fill=TEAL_BG, border=TEAL, bw=1.0)
txt(s, "Time to complete: ~30 minutes. Then billing is 100% live.", 8.62, 6.93, 4.5, 0.2, size=9, bold=True, color=TEAL)

slide_num(s, 4)

# ===================================================================
# SLIDE 5 - WHAT IS MCP? (Plain English Explanation)
# ===================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
hdr(s, "What is MCP? Why is it Phase 2?", "Model Context Protocol -- NOT a failure, it's the biggest future opportunity", accent=PURPLE)

txt(s, "MCP = Model Context Protocol (made by Anthropic, the company that makes Claude)",
    0.45, 1.52, 12.43, 0.35, size=12, color=NAVY, bold=True)
txt(s, "It's a standard that lets Claude AI connect to external tools -- Jira, GitHub, Slack, AWS CLI -- and take actions, not just suggest them.",
    0.45, 1.88, 12.43, 0.32, size=11, color=SUBTEXT)

rect(s, 0.45, 2.28, 12.43, 0.04, BORDER)

# Side by side: Today vs Phase 2
txt(s, "TODAY (Phase 1) -- Claude TELLS you what to do", 0.45, 2.42, 6.14, 0.35, size=12, bold=True, color=NAVYMD)
txt(s, "PHASE 2 (MCP) -- Claude DOES it after you approve", 6.74, 2.42, 6.14, 0.35, size=12, bold=True, color=PURPLE)

today_flow = [
    "1. VigiliCloud scans your AWS account",
    "2. Finds: S3 bucket 'my-data' is PUBLIC",
    "3. Claude explains: 'This exposes all files to internet'",
    "4. Claude shows: Run this CLI command to fix it",
    "5. YOU manually copy-paste and run the command",
    "6. YOU manually re-scan to verify it's fixed",
    "",
    "Result: 6 steps, YOU do the actual work",
]
phase2_flow = [
    "1. VigiliCloud scans your AWS account",
    "2. Finds: S3 bucket 'my-data' is PUBLIC",
    "3. Claude explains the risk AND proposes the fix",
    "4. You see: [APPROVE FIX] button -- one click",
    "5. Claude executes the fix via MCP (with least privilege)",
    "6. Auto re-scan confirms fix, audit log records everything",
    "",
    "Result: 4 steps, YOU just approve -- Claude does the work",
]
for i, (t_step, p_step) in enumerate(zip(today_flow, phase2_flow)):
    y = 2.88 + i * 0.48
    col_t = NAVYMD if "YOU" in t_step else SUBTEXT
    col_p = PURPLE if ("approve" in p_step.lower() or "Claude executes" in p_step) else SUBTEXT
    bold_t = "YOU" in t_step
    bold_p = "approve" in p_step.lower() or "Claude executes" in p_step
    if t_step:
        txt(s, t_step, 0.55, y, 6.04, 0.42, size=10.5, color=col_t, bold=bold_t)
    if p_step:
        txt(s, p_step, 6.84, y, 6.04, 0.42, size=10.5, color=col_p, bold=bold_p)

rect(s, 6.59, 2.42, 0.04, 4.5, BORDER)

# Key rule
card(s, 0.45, 6.72, 12.43, 0.62, fill=PURP_BG, border=PURPLE, bw=1.5)
txt(s, "The non-negotiable rule: Read runs automatically. Any write/fix action ALWAYS needs your approval first. Claude proposes, you decide, Claude executes.",
    0.65, 6.8, 12.0, 0.26, size=11, bold=True, color=PURPLE)
txt(s, "MCP is NOT 'AI doing things without permission' -- it's 'AI doing things with exactly one click of your approval, then doing it perfectly every time.'",
    0.65, 7.06, 12.0, 0.2, size=10, color=SUBTEXT, italic=True)

slide_num(s, 5)

# ===================================================================
# SLIDE 6 - NEXT FEATURES TO BUILD (Q2-Q3 2026)
# ===================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
hdr(s, "Next Features to Build -- Q2 & Q3 2026", "Ranked by impact. Highest value first.")

features_next = [
    ("NEXT",       "Razorpay Env Vars on Render",
     "Add 6 env vars to Render backend service. Billing goes live with zero code changes needed.",
     "30 min", "No code"),
    ("NEXT",       "Multi-Account Dashboard",
     "Single screen showing all connected accounts: risk score, last scan time, CRITICAL count, fix progress per account.",
     "2-3 days", "UI + API"),
    ("NEXT",       "Scheduled Daily Auto-Scans",
     "Cron job scans all accounts daily. Email alert only if NEW critical findings appear since last scan.",
     "1-2 days", "Backend"),
    ("PLANNED",    "Slack / Teams Webhook Alerts",
     "User pastes their Slack webhook URL in Settings. CRITICAL findings posted to their #security channel automatically.",
     "2-3 days", "Backend + UI"),
    ("PLANNED",    "SOC2 Questionnaire Autofill",
     "AI reads scan evidence and pre-fills standard SOC2 questionnaire answers. Huge value for startup enterprise deals.",
     "3-5 days", "AI + UI"),
    ("PLANNED",    "IaC Fix Snippets (Terraform / CDK)",
     "AI generates Terraform or CDK code to fix each finding -- not just CLI commands. Much higher value for DevOps teams.",
     "3-5 days", "AI + UI"),
    ("PLANNED",    "Drift Monitoring & Change Timeline",
     "Track config changes between scans. Show when each finding first appeared and what changed. Timeline view per account.",
     "4-6 days", "Backend + UI"),
    ("PLANNED",    "Developer API (REST + API Keys)",
     "API keys for developers to trigger scans and pull findings programmatically. GitHub Actions integration possible.",
     "3-4 days", "Backend"),
]

for i, (status, name, desc, effort, area) in enumerate(features_next):
    y = 1.55 + i * 0.72
    col, bg_c = STATUS_COLORS[status]
    card(s, 0.45, y, 12.43, 0.66, fill=bg_c, border=col, bw=0.75)
    rect(s, 0.45, y, 0.06, 0.66, col)
    sb = rect(s, 0.58, y + 0.18, 1.12, 0.28, col)
    sb.line.fill.background()
    txt(s, status, 0.58, y + 0.2, 1.12, 0.24, size=7.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, name, 1.82, y + 0.04, 3.2, 0.3, size=11, bold=True, color=TEXT)
    txt(s, desc, 1.82, y + 0.36, 8.4, 0.26, size=9.5, color=SUBTEXT)
    txt(s, effort, 10.38, y + 0.04, 1.32, 0.28, size=10, bold=True, color=col, align=PP_ALIGN.RIGHT)
    txt(s, area, 10.38, y + 0.36, 1.32, 0.26, size=8.5, color=MUTED, align=PP_ALIGN.RIGHT)

slide_num(s, 6)

# ===================================================================
# SLIDE 7 - Q4 2026: APPROVAL GATES + MSP
# ===================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
hdr(s, "Q4 2026 -- Approval Gates & MSP Multi-Tenant", "These two features significantly raise the product ceiling")

txt(s, "Ticket Creation & Approval Gates", 0.45, 1.52, 5.9, 0.35, size=13, bold=True, color=NAVY)
txt(s, "Planned  |  Effort: 5-8 days  |  Unlocks: enterprise workflow", 0.45, 1.88, 5.9, 0.28, size=9.5, color=AMBER)

approval_items = [
    ("Jira Ticket Auto-Creation",
     "Finding detected -> Jira issue auto-created with: severity, resource ID, affected service, and fix steps pre-filled. Engineer just needs to click resolve."),
    ("GitHub PR Draft for IaC Fixes",
     "When a Terraform/CDK fix is generated, VigiliCloud opens a PR on your repo with the exact code changes. Engineer reviews and merges."),
    ("Fix Approval Workflow",
     "Mark any AI-proposed fix as 'Needs Approval'. Assign to a team member. They get an email with one-click approve/reject. Full audit trail."),
    ("Immutable Audit Log",
     "Every scan, fix, approval, export logged with: timestamp, user, what changed, before/after values. Cannot be edited or deleted. SOC2 requirement."),
    ("Auto-Verify After Fix",
     "After a fix is marked complete, VigiliCloud auto-triggers a re-scan of just that resource to confirm the finding is resolved. Closes the loop."),
]
for i, (name, desc) in enumerate(approval_items):
    y = 2.25 + i * 0.98
    card(s, 0.45, y, 5.9, 0.86, fill=NAVY_BG, border=NAVYMD, bw=1.0)
    rect(s, 0.45, y, 0.06, 0.86, NAVYMD)
    txt(s, name, 0.63, y + 0.06, 5.55, 0.28, size=11, bold=True, color=NAVYMD)
    txt(s, desc, 0.63, y + 0.38, 5.55, 0.42, size=9.5, color=SUBTEXT)

txt(s, "MSP Multi-Tenant Dashboard", 6.6, 1.52, 6.28, 0.35, size=13, bold=True, color=NAVY)
txt(s, "Planned  |  Effort: 1-2 weeks  |  Unlocks: $999/mo MSP plan at scale", 6.6, 1.88, 6.28, 0.28, size=9.5, color=AMBER)

msp_items = [
    ("Client Account Groups",
     "Group accounts by client name. Each client has isolated accounts, scans, findings, and reports."),
    ("MSP Overview Dashboard",
     "One screen: all clients, worst risk score, last scan date, CRITICAL count, fix progress % per client."),
    ("Per-Client Report Generator",
     "One-click branded PDF/CSV report per client. Ready to attach to a monthly email to your client."),
    ("User Roles (Admin vs. Viewer)",
     "Admin = full access. Viewer = read-only. Clients can optionally be given login to see their own results."),
    ("White-Label Reports",
     "Replace 'VigiliCloud' with MSP's own logo on all reports. Agencies resell the service under their brand."),
    ("Usage Tracking Per Client",
     "Track scans and accounts per client for MSP billing. Know exactly which client uses what."),
]
for i, (name, desc) in enumerate(msp_items):
    y = 2.25 + i * 0.85
    card(s, 6.6, y, 6.28, 0.74, fill=PURP_BG, border=PURPLE, bw=0.75)
    rect(s, 6.6, y, 0.06, 0.74, PURPLE)
    txt(s, name, 6.78, y + 0.06, 5.9, 0.28, size=11, bold=True, color=PURPLE)
    txt(s, desc, 6.78, y + 0.38, 5.9, 0.3, size=9.5, color=SUBTEXT)

slide_num(s, 7)

# ===================================================================
# SLIDE 8 - PHASE 2: MCP GOVERNANCE DEEP DIVE
# ===================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
hdr(s, "Phase 2 -- MCP Governance (2027+)", "Claude doesn't just find problems -- it fixes them (with your approval)", accent=PURPLE)

txt(s, "The full MCP execution flow (Phase 2):", 0.45, 1.52, 8.0, 0.35, size=13, bold=True, color=NAVY)

flow_items = [
    ("Scan\n(reads AWS)", TEAL,   0.45),
    ("Finding\nDetected",  NAVYMD, 2.62),
    ("AI Proposes\nExact Fix", PURPLE, 4.79),
    ("You Click\nAPPROVE", AMBER,  6.96),
    ("MCP Executes\nwith Least Priv.", PURPLE, 9.13),
    ("Auto-Verify\n+ Audit Log", TEAL,   11.3),
]
for title, col, x in flow_items:
    b = rect(s, x, 2.0, 1.88, 0.88, col)
    b.line.fill.background()
    txt(s, title, x, 2.0, 1.88, 0.88, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if x < 11.3:
        txt(s, "->", x + 1.94, 2.26, 0.28, 0.35, size=14, bold=True, color=MUTED, align=PP_ALIGN.CENTER)

txt(s, "What this enables that Phase 1 cannot do:", 0.45, 3.12, 12.43, 0.32, size=12, bold=True, color=NAVY)

mcp_actions = [
    ("Fix S3 bucket in 1 click",
     "Claude applies Block Public Access to the bucket automatically. You just approve. No CLI needed.",
     TEAL),
    ("Create Jira tickets automatically",
     "Finding detected -> Jira issue created with full details -> assigned to right person -> tracked to resolution.",
     NAVYMD),
    ("Open GitHub PRs with code fixes",
     "Terraform/CDK code generated -> PR opened on your repo -> engineer reviews and merges -> issue resolved.",
     NAVYMD),
    ("Slack approval from your phone",
     "CRITICAL finding -> Slack message with [APPROVE FIX] button -> approve from Slack -> Claude fixes it.",
     PURPLE),
    ("Multi-cloud (Azure + GCP)",
     "Same approval flow, same audit log -- but fixing misconfigs across Azure and GCP accounts too.",
     PURPLE),
    ("SOC2 questionnaire answered",
     "Enterprise customer sends security questionnaire -> Claude fills all answers backed by scan evidence -> you review and send.",
     AMBER),
]
for i, (action, desc, col) in enumerate(mcp_actions):
    ci = i % 2; row = i // 2
    x = 0.45 + ci * 6.44
    y = 3.55 + row * 1.1
    card(s, x, y, 6.14, 0.98, fill=PURP_BG, border=col, bw=1.0)
    rect(s, x, y, 0.06, 0.98, col)
    txt(s, action, x + 0.2, y + 0.06, 2.8, 0.3, size=11, bold=True, color=col)
    txt(s, desc, x + 0.2, y + 0.44, 5.74, 0.46, size=10, color=SUBTEXT)

card(s, 0.45, 6.72, 12.43, 0.62, fill=NAVY_BG, border=NAVY, bw=1.5)
txt(s, "Why this creates an unbeatable moat: Competitors give you alerts. VigiliCloud Phase 2 closes the full loop from detection to verified fix -- with one click.",
    0.65, 6.8, 12.0, 0.26, size=11, bold=True, color=NAVY)
txt(s, "MCP is already used by thousands of Claude integrations. We reuse Anthropic's standard -- zero infrastructure to build from scratch.",
    0.65, 7.06, 12.0, 0.22, size=10, color=SUBTEXT, italic=True)

slide_num(s, 8)

# ===================================================================
# SLIDE 9 - MASTER TIMELINE
# ===================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
hdr(s, "Master Roadmap Timeline -- All Features", "From today's MVP to Phase 2 MCP platform")

timeline = [
    ("NOW\nMay 2026",    TEAL,   [
        ("DONE",    "Core scanner (10 checks)", TEAL),
        ("DONE",    "Claude AI analysis",        TEAL),
        ("DONE",    "Exports + Email alerts",    TEAL),
        ("NEXT",    "Razorpay env vars (30min)", AMBER),
    ]),
    ("Q3 2026\nJul-Sep", NAVYMD, [
        ("PLANNED", "Multi-account dashboard",   NAVYMD),
        ("PLANNED", "Scheduled daily scans",     NAVYMD),
        ("PLANNED", "Slack/Teams alerts",         NAVYMD),
        ("PLANNED", "SOC2 questionnaire AI",     NAVYMD),
    ]),
    ("Q4 2026\nOct-Dec", ORANGE, [
        ("PLANNED", "Jira + GitHub PR drafts",   ORANGE),
        ("PLANNED", "Approval gates + audit log",ORANGE),
        ("PLANNED", "MSP multi-tenant mode",     ORANGE),
        ("PLANNED", "Developer API access",      ORANGE),
    ]),
    ("2027+\nPhase 2",   PURPLE, [
        ("FUTURE",  "MCP execution engine",      PURPLE),
        ("FUTURE",  "Azure + GCP support",       PURPLE),
        ("FUTURE",  "White-label MSP mode",      PURPLE),
        ("FUTURE",  "HIPAA / PCI custom checks", PURPLE),
    ]),
]

rect(s, 0.45, 2.24, 12.43, 0.06, BORDER)

for qi, (quarter, col, items) in enumerate(timeline):
    x = 0.45 + qi * 3.22
    rect(s, x + 1.41, 2.12, 0.26, 0.26, col)
    txt(s, quarter, x, 1.52, 3.1, 0.55, size=11, bold=True, color=col, align=PP_ALIGN.CENTER)
    card(s, x, 2.44, 3.05, 4.68, fill=WHITE, border=col, bw=1.5)
    rect(s, x, 2.44, 3.05, 0.08, col)

    for j, (status, task, task_col) in enumerate(items):
        ty = 2.6 + j * 1.08
        stat_col, stat_bg = STATUS_COLORS.get(status, (SLATE, SLT_BG))
        card(s, x + 0.14, ty, 2.77, 0.9, fill=stat_bg, border=stat_col, bw=0.5)
        sb = rect(s, x + 0.14, ty + 0.04, 0.96, 0.24, stat_col)
        sb.line.fill.background()
        txt(s, status, x + 0.14, ty + 0.05, 0.96, 0.22, size=6.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, task, x + 0.24, ty + 0.36, 2.54, 0.46, size=10, color=TEXT)

txt(s, "Timelines are estimates. Features may shift based on customer feedback and which features drive most revenue.",
    0.45, 7.22, 12.43, 0.22, size=9.5, color=MUTED, italic=True, align=PP_ALIGN.CENTER)

slide_num(s, 9)

# ===================================================================
# SLIDE 10 - IMMEDIATE 30-DAY ACTION PLAN
# ===================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
hdr(s, "Immediate Action Plan -- Next 30 Days", "The highest-leverage tasks right now, in priority order")

weeks = [
    ("This Week", AMBER, [
        ("Add Razorpay env vars to Render",             "30 min",  "Billing live instantly"),
        ("Test full payment flow end-to-end",           "2 hrs",   "Verify subscribe + cancel works"),
        ("Set up vigilicloud.com email (Resend DNS)",   "2 hrs",   "Replace onboarding@resend.dev"),
        ("LinkedIn: 10 AWS consultant outreaches/day",  "Daily",   "First paying customers"),
    ]),
    ("Week 2", NAVYMD, [
        ("Build multi-account dashboard UI",            "2-3 days","Required for Pro/MSP plan value"),
        ("Add scheduled scan cron job (backend)",       "1-2 days","Automates daily scanning"),
        ("Create ProductHunt launch page + screenshots","1 day",   "Prepare for coordinated launch"),
        ("Build list of 30 target MSPs + SOC2 startups","2 hrs",  "GTM foundation"),
    ]),
    ("Week 3-4", TEAL, [
        ("ProductHunt + Hacker News launch",            "Launch day","Target: 1,000+ signups"),
        ("Add Slack webhook integration (Settings)",    "2-3 days","Most common early ask"),
        ("Cold email campaign to 30 target accounts",  "Ongoing", "High-intent outreach"),
        ("Close first 3 paying pilots",                "Ongoing", "Decision gate: proves PMF"),
    ]),
]

for wi, (week, col, tasks) in enumerate(weeks):
    x = 0.45 + wi * 4.3
    rect(s, x, 1.52, 4.05, 0.42, col)
    txt(s, week, x, 1.52, 4.05, 0.42, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    card(s, x, 1.94, 4.05, 5.22, fill=WHITE, border=col, bw=1.0)
    for ti, (task, effort, result) in enumerate(tasks):
        ty = 2.02 + ti * 1.2
        card(s, x + 0.12, ty, 3.81, 1.05, fill=LGRAY, border=BORDER, bw=0.5)
        txt(s, task, x + 0.25, ty + 0.06, 3.55, 0.42, size=10, bold=True, color=TEXT)
        effort_b = rect(s, x + 0.25, ty + 0.56, 0.88, 0.24, col)
        effort_b.line.fill.background()
        txt(s, effort, x + 0.25, ty + 0.57, 0.88, 0.22, size=7.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, result, x + 1.22, ty + 0.58, 2.68, 0.22, size=9.5, color=SUBTEXT, italic=True)

card(s, 0.45, 7.2, 12.43, 0.22, fill=TEAL_BG, border=TEAL, bw=1.0)
txt(s, "Decision Gate: If 3 paying pilots close in 30 days -> business proven. Double down on features + GTM. Original plan from Feb 2026.",
    0.65, 7.23, 12.0, 0.18, size=10, bold=True, color=TEAL)

slide_num(s, 10)

# -- Save --
out = r"c:\Users\leela\compliance-ai-saas\VigiliCloud_Internal_Roadmap.pptx"
prs.save(out)
print(f"Saved: {out}")
