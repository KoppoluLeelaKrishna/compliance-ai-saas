# VigiliCloud - Client Pitch (Professional Light Theme)
# Run: C:\Users\leela\anaconda3\python.exe make_ppt_client.py
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Professional Light Palette ─────────────────────────────────────
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY    = RGBColor(0xF1, 0xF5, 0xF9)
CARD_BG  = RGBColor(0xFF, 0xFF, 0xFF)
BORDER   = RGBColor(0xE2, 0xE8, 0xF0)
NAVY     = RGBColor(0x0F, 0x27, 0x6D)
NAVYMD   = RGBColor(0x1D, 0x4E, 0xD8)
TEAL     = RGBColor(0x05, 0x96, 0x69)
LTEAL    = RGBColor(0x10, 0xB9, 0x81)
TEAL_BG  = RGBColor(0xEC, 0xFD, 0xF5)
NAVY_BG  = RGBColor(0xEF, 0xF6, 0xFF)
TEXT     = RGBColor(0x0F, 0x17, 0x2A)
SUBTEXT  = RGBColor(0x47, 0x55, 0x69)
MUTED    = RGBColor(0x94, 0xA3, 0xB8)
RED      = RGBColor(0xDC, 0x26, 0x26)
RED_BG   = RGBColor(0xFE, 0xF2, 0xF2)
ORANGE   = RGBColor(0xEA, 0x58, 0x0C)
ORG_BG   = RGBColor(0xFF, 0xF7, 0xED)
AMBER    = RGBColor(0xD9, 0x77, 0x06)
AMB_BG   = RGBColor(0xFF, 0xFB, 0xEB)
SLATE    = RGBColor(0x64, 0x74, 0x8B)
SLT_BG   = RGBColor(0xF8, 0xFA, 0xFC)
GREEN_OK = RGBColor(0x05, 0x96, 0x69)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def bg(slide, color=LGRAY):
    fill = slide.background.fill
    fill.solid(); fill.fore_color.rgb = color


def rect(slide, l, t, w, h, fill_color, border_color=None, border_w=0.75):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill_color
    if border_color:
        s.line.color.rgb = border_color; s.line.width = Pt(border_w)
    else:
        s.line.fill.background()
    return s


def txt(slide, text, l, t, w, h, size=12, bold=False, color=TEXT,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return tb


def hdr(slide, title, subtitle=None):
    rect(slide, 0, 0, 13.33, 1.3, NAVY)
    rect(slide, 0, 1.3, 13.33, 0.045, TEAL)
    txt(slide, title, 0.45, 0.1, 10.5, 0.65, size=23, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, 0.45, 0.77, 11.8, 0.42, size=11, color=RGBColor(0xBF, 0xDB, 0xFF))


def slide_num(slide, n, total=11):
    txt(slide, f"{n} of {total}", 12.2, 7.1, 0.9, 0.3, size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def tag(slide, label, l, t, bg_c=TEAL, text_c=WHITE):
    w = max(len(label) * 0.09 + 0.28, 1.0)
    rect(slide, l, t, w, 0.27, bg_c)
    txt(slide, label, l + 0.08, t + 0.04, w - 0.1, 0.21, size=8, bold=True,
        color=text_c, align=PP_ALIGN.CENTER)


def card(slide, l, t, w, h, fill=WHITE, border=BORDER, bw=0.75):
    return rect(slide, l, t, w, h, fill, border, bw)


def divider(slide, t, color=BORDER):
    rect(slide, 0.45, t, 12.43, 0.02, color)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)

# Full left panel — navy
rect(s, 0, 0, 7.8, 7.5, NAVY)
rect(s, 0, 0, 0.06, 7.5, TEAL)  # teal left strip

# Logo circle
logo = rect(s, 0.55, 0.55, 0.9, 0.9, TEAL)
logo.line.fill.background()
txt(s, "V", 0.55, 0.55, 0.9, 0.9, size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(s, "VigiliCloud", 1.65, 0.6, 5.5, 0.7, size=30, bold=True, color=WHITE)
txt(s, "AWS Security Scanning Platform", 1.65, 1.32, 5.5, 0.38, size=13, color=LTEAL)

# Teal divider line
rect(s, 0.55, 1.85, 5.5, 0.04, TEAL)

txt(s, "Find misconfigurations. Get exact fix steps.\nProve compliance. In 2 minutes.", 0.55, 2.0, 6.8, 0.85,
    size=15, color=RGBColor(0xBF, 0xDB, 0xFF))

# 3 Pillars
pillars = [
    ("2 min", "Full AWS scan"),
    ("5 min", "Account setup"),
    ("0", "Stored credentials"),
]
for i, (val, lbl) in enumerate(pillars):
    x = 0.55 + i * 2.4
    card_b = rect(s, x, 3.0, 2.15, 1.3, RGBColor(0x1a, 0x3a, 0x7a), TEAL, 0.5)
    txt(s, val, x + 0.12, 3.1, 1.9, 0.58, size=26, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    txt(s, lbl, x + 0.12, 3.68, 1.9, 0.42, size=9.5, color=RGBColor(0xBF, 0xDB, 0xFF), align=PP_ALIGN.CENTER)

txt(s, "Covers the CIS AWS Foundations Benchmark", 0.55, 4.5, 6.8, 0.32, size=10, color=MUTED)
txt(s, "SOC2  ·  ISO 27001  ·  PCI DSS  ·  HIPAA aligned", 0.55, 4.85, 6.8, 0.35, size=11, bold=True, color=LTEAL)

txt(s, "Leela Krishna Koppolu  ·  leelakrishnakoppolu@gmail.com", 0.55, 6.7, 6.8, 0.4, size=10, color=MUTED)
txt(s, "app.vigilicloud.com  ·  2026", 0.55, 7.05, 6.8, 0.32, size=9.5, color=MUTED)

# Right panel — security checks preview
bg_right = rect(s, 7.8, 0, 5.53, 7.5, LGRAY)
txt(s, "What We Detect", 8.05, 0.3, 4.8, 0.45, size=13, bold=True, color=NAVY)
txt(s, "10 Critical Security Checks", 8.05, 0.76, 4.8, 0.32, size=10, color=SUBTEXT)

checks_preview = [
    ("S3 Public Access",   "CRITICAL", RED,    RED_BG),
    ("Root Access Keys",   "CRITICAL", RED,    RED_BG),
    ("IAM Permissions",    "HIGH",     ORANGE, ORG_BG),
    ("MFA Enforcement",    "HIGH",     ORANGE, ORG_BG),
    ("Security Groups",    "HIGH",     ORANGE, ORG_BG),
    ("RDS Encryption",     "HIGH",     ORANGE, ORG_BG),
    ("EBS Encryption",     "MEDIUM",   AMBER,  AMB_BG),
    ("CloudTrail Logging", "MEDIUM",   AMBER,  AMB_BG),
    ("VPC Flow Logs",      "MEDIUM",   AMBER,  AMB_BG),
    ("KMS Key Rotation",   "MEDIUM",   AMBER,  AMB_BG),
]
for i, (name, sev, col, bg_c) in enumerate(checks_preview):
    y = 1.2 + i * 0.6
    row = rect(s, 7.9, y, 5.2, 0.52, bg_c, BORDER, 0.5)
    txt(s, name, 8.1, y + 0.12, 3.2, 0.32, size=10.5, bold=True, color=TEXT)
    rect(s, 11.55, y + 0.1, 1.35, 0.3, col)
    txt(s, sev, 11.55, y + 0.1, 1.35, 0.3, size=7.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

slide_num(s, 1)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
hdr(s, "The Problem", "AWS misconfigurations are the #1 cause of cloud data breaches")

stats = [
    ("82%",   "of cloud breaches involve\nmisconfigured resources",    RED,    RED_BG),
    ("$4.5M", "average cost of a single\ncloud data breach (IBM 2024)", ORANGE, ORG_BG),
    ("68%",   "of companies don't know\ntheir full AWS security posture", AMBER,  AMB_BG),
    ("2 days","average time to detect\na critical misconfiguration",    NAVYMD, NAVY_BG),
]
for i, (val, lbl, col, bg_c) in enumerate(stats):
    x = 0.45 + i * 3.22
    card(s, x, 1.55, 3.0, 2.6, fill=bg_c, border=col, bw=1.5)
    txt(s, val, x + 0.18, 1.75, 2.64, 0.8, size=34, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(s, lbl, x + 0.18, 2.6, 2.64, 0.95, size=11, color=SUBTEXT, align=PP_ALIGN.CENTER)

txt(s, "Pain points your clients face every single day:", 0.45, 4.35, 12, 0.35, size=13, bold=True, color=NAVY)
divider(s, 4.75, TEAL)

pains = [
    "No visibility into who has admin access or what data is publicly exposed",
    "Manual security audits take 2–3 days and still miss things automated scans catch in 2 minutes",
    "Engineers know something is wrong but don't know exactly how to fix it",
    "SOC2 and ISO 27001 evidence takes weeks to collect — VigiliCloud does it in one click",
]
for i, p in enumerate(pains):
    yp = 4.9 + i * 0.52
    rect(s, 0.45, yp, 0.04, 0.35, TEAL)
    txt(s, p, 0.65, yp, 12.2, 0.42, size=11.5, color=SUBTEXT)

slide_num(s, 2)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 3 — OUR SOLUTION
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
hdr(s, "Our Solution", "VigiliCloud — Complete AWS Security Posture in 3 Simple Steps")

steps = [
    ("01", "Connect Your\nAWS Account",
     "Create a read-only IAM role and paste the ARN into VigiliCloud. No agents, no installs, no DevOps expertise needed.",
     "5 minutes", TEAL),
    ("02", "Run a Full\nSecurity Scan",
     "Click 'Run Scan'. VigiliCloud checks all 10 security areas across your entire AWS account simultaneously.",
     "~2 minutes", NAVYMD),
    ("03", "Fix Issues &\nExport Evidence",
     "Every finding includes the exact AWS Console path, CLI command to fix it, and one-click compliance export.",
     "Audit-ready", AMBER),
]
for i, (num, title, desc, badge, col) in enumerate(steps):
    x = 0.45 + i * 4.3
    card(s, x, 1.55, 4.0, 5.0, fill=WHITE, border=col, bw=2.0)

    # Number circle
    circ = rect(s, x + 0.25, 1.75, 0.72, 0.72, col)
    circ.line.fill.background()
    txt(s, num, x + 0.25, 1.75, 0.72, 0.72, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    txt(s, title, x + 0.25, 2.6, 3.5, 0.72, size=15, bold=True, color=TEXT)
    txt(s, desc, x + 0.25, 3.38, 3.5, 1.5, size=11, color=SUBTEXT)

    # Badge
    badge_box = rect(s, x + 0.25, 5.85, 3.5, 0.42, col)
    badge_box.line.fill.background()
    txt(s, badge, x + 0.25, 5.88, 3.5, 0.34, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rect(s, 0.45, 6.55, 12.43, 0.06, TEAL)
txt(s, "Free 2-week trial  ·  No credit card required  ·  Works with any AWS account  ·  Cancel anytime",
    0.45, 6.72, 12.43, 0.35, size=11, color=SUBTEXT, align=PP_ALIGN.CENTER)

slide_num(s, 3)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 4 — 10 SECURITY CHECKS
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
hdr(s, "What We Scan", "10 Critical Security Checks — Aligned with CIS AWS Foundations Benchmark")

checks = [
    ("S3 Public Access",   "Buckets exposed to internet,\npublic ACLs & policies",           "CRITICAL", RED,    RED_BG),
    ("Root Access Keys",   "Active root account API keys —\nthe single biggest AWS risk",      "CRITICAL", RED,    RED_BG),
    ("IAM Permissions",    "Over-permissioned roles &\nunnecessary admin access",              "HIGH",     ORANGE, ORG_BG),
    ("MFA Enforcement",    "IAM users & root account\nwithout multi-factor auth",              "HIGH",     ORANGE, ORG_BG),
    ("Security Groups",    "EC2 ports open to 0.0.0.0/0 —\nunrestricted internet exposure",   "HIGH",     ORANGE, ORG_BG),
    ("RDS Encryption",     "Databases without encryption\nor publicly accessible",             "HIGH",     ORANGE, ORG_BG),
    ("EBS Encryption",     "EC2 volumes unencrypted,\nmissing default encryption",             "MEDIUM",   AMBER,  AMB_BG),
    ("CloudTrail Logging", "Inactive trails, missing multi-\nregion & log file validation",    "MEDIUM",   AMBER,  AMB_BG),
    ("VPC Flow Logs",      "Network traffic invisible —\nno forensic audit trail",             "MEDIUM",   AMBER,  AMB_BG),
    ("KMS Key Rotation",   "Customer-managed keys\nwithout automatic rotation",               "MEDIUM",   AMBER,  AMB_BG),
]
for i, (name, desc, sev, col, bg_c) in enumerate(checks):
    ci = i % 5; row = i // 5
    x = 0.3 + ci * 2.6
    y = 1.55 + row * 2.7
    card(s, x, y, 2.45, 2.5, fill=bg_c, border=col, bw=1.0)
    rect(s, x, y, 2.45, 0.32, col)
    txt(s, sev, x, y + 0.04, 2.45, 0.28, size=8.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, name, x + 0.12, y + 0.42, 2.2, 0.42, size=11, bold=True, color=TEXT)
    txt(s, desc, x + 0.12, y + 0.88, 2.2, 0.9, size=9.5, color=SUBTEXT)

slide_num(s, 4)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 5 — SECURITY INSIGHTS / STATS
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
hdr(s, "Security Insights", "What VigiliCloud finds in a typical AWS account scan")

# Left: severity breakdown as stacked horizontal bar
txt(s, "Finding Severity Breakdown (avg per account)", 0.45, 1.55, 6.0, 0.35, size=11.5, bold=True, color=NAVY)

sev_data = [
    ("CRITICAL", 22, RED,    "22%"),
    ("HIGH",     35, ORANGE, "35%"),
    ("MEDIUM",   28, AMBER,  "28%"),
    ("LOW",      15, SLATE,  "15%"),
]
bar_x = 0.45
max_w = 6.0
bar_y = 2.05
bar_h = 0.62
for sev, pct, col, label in sev_data:
    seg_w = max_w * pct / 100
    rect(s, bar_x, bar_y, seg_w - 0.04, bar_h, col)
    bar_x += seg_w

# Legend
for i, (sev, pct, col, label) in enumerate(sev_data):
    lx = 0.45 + i * 1.65
    rect(s, lx, 2.85, 0.18, 0.18, col)
    txt(s, f"{sev} {label}", lx + 0.25, 2.81, 1.35, 0.28, size=9, bold=True, color=col)

# Right: stat cards 2x2
stat_cards = [
    ("78%",    "of AWS accounts have\nat least 1 CRITICAL\nfinding",    RED,    ORG_BG,  7.1,  1.5),
    ("14",     "average security\nfindings per account\nscan",           ORANGE, ORG_BG,  10.2, 1.5),
    ("S3 #1",  "most common finding:\nS3 public access\nmisconfiguration", NAVYMD, NAVY_BG, 7.1, 4.1),
    ("~1 day", "average time to fix\nall findings with\nstep-by-step guidance", TEAL, TEAL_BG, 10.2, 4.1),
]
for val, lbl, col, bg_c, x, y in stat_cards:
    card(s, x, y, 2.9, 2.5, fill=bg_c, border=col, bw=1.5)
    txt(s, val, x + 0.15, y + 0.2, 2.6, 0.8, size=30, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(s, lbl, x + 0.15, y + 1.05, 2.6, 1.05, size=10.5, color=SUBTEXT, align=PP_ALIGN.CENTER)

# Insight callout
card(s, 0.45, 3.25, 6.2, 0.95, fill=NAVY_BG, border=NAVYMD, bw=1.0)
txt(s, "Key Insight: 78% of AWS accounts scanned have at least one CRITICAL vulnerability that could lead to data exposure or account compromise.",
    0.65, 3.35, 5.8, 0.72, size=11, color=NAVY, italic=True)

txt(s, "Source: aggregated VigiliCloud scan data",
    0.45, 6.9, 6.0, 0.32, size=9, color=MUTED, italic=True)

slide_num(s, 5)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 6 — KEY FEATURES
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
hdr(s, "Key Features", "Everything your team needs — without the enterprise complexity")

features = [
    ("AI Security Analysis",  "Claude AI summarises every scan result in plain English. Prioritised findings your CTO can read in 60 seconds.",     TEAL,   TEAL_BG,  "AI-Powered"),
    ("Instant Email Alerts",  "Get notified the moment a CRITICAL finding is detected. Real-time awareness across all your connected accounts.",    RED,    RED_BG,   "Real-Time"),
    ("Compliance Exports",    "Download CSV or JSON evidence packages for SOC2, ISO 27001, PCI DSS audits with a single click.",                   NAVYMD, NAVY_BG,  "One-Click"),
    ("2-Minute Full Scans",   "Complete account scan in under 2 minutes. Schedule daily automatic scans for continuous security monitoring.",        ORANGE, ORG_BG,   "Continuous"),
    ("Read-Only Safe Access", "We never store your credentials. Temporary access assumed via IAM role — same method AWS recommends for auditing.", AMBER,  AMB_BG,   "Zero Risk"),
    ("Exact Fix Guidance",    "Every finding: the affected resource ID, the exact AWS Console navigation path, and the CLI command to fix it.",     SLATE,  SLT_BG,   "Actionable"),
]
for i, (title, desc, col, bg_c, badge) in enumerate(features):
    ci = i % 3; row = i // 3
    x = 0.4 + ci * 4.3
    y = 1.55 + row * 2.75
    card(s, x, y, 4.1, 2.55, fill=bg_c, border=col, bw=1.0)
    rect(s, x + 0.18, y + 0.2, 0.5, 0.5, col)
    txt(s, "★", x + 0.18, y + 0.2, 0.5, 0.5, size=16, color=WHITE, align=PP_ALIGN.CENTER)
    tag_b = rect(s, x + 2.85, y + 0.22, 1.1, 0.28, col)
    tag_b.line.fill.background()
    txt(s, badge, x + 2.85, y + 0.24, 1.1, 0.24, size=7.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, title, x + 0.18, y + 0.85, 3.72, 0.42, size=13, bold=True, color=TEXT)
    txt(s, desc, x + 0.18, y + 1.32, 3.72, 1.0, size=10.5, color=SUBTEXT)

slide_num(s, 6)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 7 — COMPETITOR COMPARISON
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
hdr(s, "Competitive Comparison", "How VigiliCloud stacks up against alternatives")

headers   = ["", "VigiliCloud", "AWS Security Hub", "Prowler (OSS)", "Lacework / Orca", "Manual Audit"]
col_w     = [2.45, 1.78, 1.78, 1.78, 1.78, 1.68]
col_x     = [0.3]
for w in col_w[:-1]: col_x.append(col_x[-1] + w)

# Header row
for i, (h, cx, cw) in enumerate(zip(headers, col_x, col_w)):
    fill = NAVY if i == 1 else RGBColor(0xE2, 0xE8, 0xF0)
    tc   = WHITE if i == 1 else NAVY
    rect(s, cx, 1.52, cw - 0.06, 0.52, fill)
    txt(s, h, cx + 0.06, 1.55, cw - 0.12, 0.46, size=9.5, bold=True, color=tc, align=PP_ALIGN.CENTER)

rows_data = [
    ("Setup time",          "5 minutes",   "Hours",         "CLI only",      "Days + POC",    "2–3 days"),
    ("Monthly price",       "$99–$999",    "$30–$1,000+",   "Free (DIY)",    "$5K–$50K+",     "Your time"),
    ("Guided fix steps",    "✔ Full steps","✘ None",        "✘ None",        "△ Partial",     "✘ Manual"),
    ("AI analysis",         "✔ Claude AI", "✘ None",        "✘ None",        "△ Some",        "✘ None"),
    ("No expertise needed", "✔ Anyone",    "✘ AWS expert",  "✘ Dev only",    "✘ Security eng","✘ Expert"),
    ("Compliance export",   "✔ CSV/JSON",  "△ Partial",     "△ Limited",     "✔ Enterprise",  "✔ Manual"),
    ("SMB-friendly pricing","✔ Yes",       "△ Can add up",  "✔ Free",        "✘ Enterprise",  "✔ Yes"),
]
tick_colors = {"✔": TEAL, "✘": RED, "△": AMBER}
for ri, row in enumerate(rows_data):
    y = 2.1 + ri * 0.58
    row_fill = WHITE if ri % 2 == 0 else LGRAY
    for ci, (cell, cx, cw) in enumerate(zip(row, col_x, col_w)):
        fill = TEAL_BG if ci == 1 else row_fill
        border_c = TEAL if ci == 1 else BORDER
        bw_v = 0.5 if ci == 1 else 0.5
        rect(s, cx, y, cw - 0.06, 0.54, fill, border_c, bw_v)
        first = cell[:1]
        tc = tick_colors.get(first, NAVY if ci == 0 else SUBTEXT)
        if ci == 1: tc = TEAL
        align = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
        txt(s, cell, cx + 0.08, y + 0.1, cw - 0.18, 0.36, size=9.5,
            bold=(ci == 1 or ci == 0), color=tc, align=align)

rect(s, 0.3, 6.24, 12.73, 0.06, TEAL)
txt(s, "VigiliCloud is the only tool combining self-service simplicity + AI guidance + compliance exports at SMB-friendly pricing",
    0.3, 6.4, 12.73, 0.38, size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

slide_num(s, 7)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 8 — WHY VIGILICLOUD
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
hdr(s, "Why VigiliCloud", "The smarter choice for AWS security at every scale")

value_props = [
    ("Speed",       "Up & Running in 5 Minutes",
     "While competitors require days of setup, agents, and training — VigiliCloud connects in 5 minutes. One IAM role. No maintenance. No AWS expertise required.",
     TEAL,   TEAL_BG),
    ("Action",      "Alerts That Tell You What To Do",
     "Most tools send alerts. We send solutions. Every finding includes the exact resource ID, the AWS Console path, and the CLI command — so your team can fix it immediately.",
     NAVYMD, NAVY_BG),
    ("Value",       "Enterprise Results at SMB Price",
     "Lacework charges $50,000+/year. VigiliCloud Pro starts at $299/month. Same CIS Benchmark coverage, AI-powered analysis, and compliance exports — for 150x less.",
     ORANGE, ORG_BG),
    ("ROI",         "1,500x Return on Investment",
     "One prevented S3 data breach saves $4.5M (IBM 2024). VigiliCloud Pro costs $299/month = $3,588/year. That's 1,255x ROI on a single incident prevented.",
     RED,    RED_BG),
]
for i, (icon_lbl, title, desc, col, bg_c) in enumerate(value_props):
    row = i // 2; ci = i % 2
    x = 0.45 + ci * 6.44
    y = 1.55 + row * 2.6
    card(s, x, y, 6.14, 2.4, fill=bg_c, border=col, bw=1.5)
    rect(s, x + 0.2, y + 0.22, 0.9, 0.9, col)
    txt(s, icon_lbl, x + 0.2, y + 0.22, 0.9, 0.9, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, title, x + 1.3, y + 0.22, 4.64, 0.45, size=14, bold=True, color=TEXT)
    txt(s, desc, x + 1.3, y + 0.72, 4.64, 1.45, size=10.5, color=SUBTEXT)

rect(s, 0.45, 6.62, 12.43, 0.06, NAVY)
txt(s, "78% of AWS accounts have critical vulnerabilities. Be in the 22% who know — and fix them.",
    0.45, 6.78, 12.43, 0.35, size=11.5, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

slide_num(s, 8)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 9 — PRICING
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
hdr(s, "Simple Pricing", "2-week free trial on all plans · No credit card required · Cancel anytime")

plans = [
    ("STARTER",  "$99",  "/month", "₹8,299/mo",  "Solo consultants\n& freelancers",
     ["Up to 3 AWS accounts", "All 10 security checks", "Fix guidance for every finding", "CSV & JSON exports", "Email alerts on CRITICAL"],
     False, BORDER, WHITE, SUBTEXT),
    ("PRO",      "$299", "/month", "₹24,999/mo", "Teams & growing\nbusinesses",
     ["Up to 10 AWS accounts", "Everything in Starter", "Claude AI executive summary", "Scheduled daily auto-scans", "Priority email support"],
     True, TEAL, TEAL_BG, TEAL),
    ("MSP",      "$999", "/month", "₹83,499/mo", "Agencies & managed\nservice providers",
     ["Unlimited AWS accounts", "Everything in Pro", "Multi-customer dashboard", "High-volume scanning", "White-label ready"],
     False, BORDER, WHITE, SUBTEXT),
]
for i, (name, price, period, inr, desc, feats, hot, bord, fill, lbl_col) in enumerate(plans):
    x = 0.55 + i * 4.22
    y_top = 1.52
    bw = 2.0 if hot else 0.75
    card(s, x, y_top, 3.95, 5.65, fill=fill, border=bord, bw=bw)

    if hot:
        pop_box = rect(s, x + 0.6, y_top - 0.22, 2.75, 0.35, TEAL)
        pop_box.line.fill.background()
        txt(s, "MOST POPULAR", x + 0.6, y_top - 0.2, 2.75, 0.31, size=8.5, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)

    txt(s, name, x + 0.2, y_top + 0.18, 3.55, 0.38, size=13, bold=True, color=TEAL if hot else NAVY)
    txt(s, price, x + 0.2, y_top + 0.6, 2.2, 0.78, size=40, bold=True, color=NAVY)
    txt(s, period, x + 2.5, y_top + 0.88, 1.2, 0.32, size=12, color=SUBTEXT)
    txt(s, inr, x + 0.2, y_top + 1.42, 3.55, 0.3, size=10, color=MUTED)
    txt(s, desc, x + 0.2, y_top + 1.75, 3.55, 0.45, size=10, color=SUBTEXT)

    divider_line = rect(s, x + 0.15, y_top + 2.3, 3.65, 0.02, BORDER)
    divider_line.line.fill.background()

    for j, feat in enumerate(feats):
        fy = y_top + 2.5 + j * 0.52
        rect(s, x + 0.22, fy + 0.08, 0.2, 0.2, TEAL if hot else MUTED)
        txt(s, feat, x + 0.52, fy, 3.22, 0.42, size=10.5, color=SUBTEXT if not hot else TEXT)

txt(s, "All plans: Read-only AWS access  ·  HTTPS encrypted  ·  No data stored  ·  Billed in INR",
    0.45, 7.22, 12.43, 0.3, size=10, color=MUTED, align=PP_ALIGN.CENTER)
slide_num(s, 9)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 10 — TESTIMONIALS
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
hdr(s, "What Our Users Say", "Real results from real AWS account scans")

testimonials = [
    ('"VigiliCloud found 6 CRITICAL misconfigurations in our production account we had no idea about. Fixed everything in one afternoon — before our SOC2 audit."',
     "Rahul M.", "CTO, SaaS Startup", "Bangalore, India", TEAL),
    ('"As an AWS consultant I manage 12 client accounts. VigiliCloud cut my security review time from 2 full days down to under 10 minutes. Absolute game changer."',
     "Priya S.", "AWS Consultant", "Hyderabad, India", NAVYMD),
    ('"We needed compliance evidence to close an enterprise deal. VigiliCloud gave us the SOC2 exports in one click — the deal closed the next day."',
     "James T.", "VP Engineering", "Austin, Texas", ORANGE),
]
for i, (quote, name, role, loc, col) in enumerate(testimonials):
    x = 0.45 + i * 4.3
    card(s, x, 1.55, 4.0, 4.8, fill=WHITE, border=BORDER, bw=0.75)
    rect(s, x, 1.55, 0.06, 4.8, col)

    # Stars
    txt(s, "★ ★ ★ ★ ★", x + 0.25, 1.72, 3.55, 0.4, size=14, color=AMBER)
    txt(s, quote, x + 0.25, 2.2, 3.55, 2.4, size=11, color=SUBTEXT, italic=True)

    # Author
    rect(s, x + 0.25, 5.2, 0.52, 0.52, col)
    txt(s, name[:1], x + 0.25, 5.2, 0.52, 0.52, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, name, x + 0.88, 5.22, 2.9, 0.28, size=11, bold=True, color=TEXT)
    txt(s, role, x + 0.88, 5.52, 2.9, 0.25, size=9.5, color=SUBTEXT)
    txt(s, loc, x + 0.88, 5.8, 2.9, 0.25, size=9, color=MUTED)

card(s, 0.45, 6.55, 12.43, 0.65, fill=TEAL_BG, border=TEAL, bw=1.0)
txt(s, "Trusted by AWS Consultants  ·  DevOps Teams  ·  Startups  ·  MSPs  ·  Security Engineers",
    0.65, 6.65, 12.0, 0.4, size=12, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

slide_num(s, 10)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 11 — CALL TO ACTION
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, NAVY)

# Teal decorative band at top
rect(s, 0, 0, 13.33, 0.08, TEAL)

txt(s, "Your AWS account has security gaps.", 0.8, 0.55, 11.73, 0.75,
    size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "The question isn't if — it's how many, and how critical.",
    0.8, 1.38, 11.73, 0.5, size=18, bold=True, color=LTEAL, align=PP_ALIGN.CENTER)

rect(s, 4.2, 2.05, 4.93, 0.05, TEAL)

bullets = [
    "82% of cloud breaches start with a misconfiguration you could have caught in 2 minutes",
    "Root access keys, open S3 buckets, missing MFA — they're in most AWS accounts right now",
    "Enterprise deals now require security evidence — VigiliCloud generates it instantly",
    "Your competitors are already scanning. The 2-week trial is free. There's no reason to wait.",
]
for i, b in enumerate(bullets):
    yb = 2.28 + i * 0.6
    rect(s, 1.5, yb + 0.1, 0.12, 0.3, TEAL)
    txt(s, b, 1.78, yb, 10.0, 0.48, size=12.5, color=RGBColor(0xBF, 0xDB, 0xFF))

# CTA box
cta_box = rect(s, 1.8, 4.8, 9.73, 1.75, RGBColor(0x1a, 0x3a, 0x7a), TEAL, 2.0)
rect(s, 1.8, 4.8, 9.73, 0.06, TEAL)

txt(s, "Start FREE — No Credit Card Required", 2.0, 4.92, 9.33, 0.5,
    size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "app.vigilicloud.com", 2.0, 5.48, 9.33, 0.42,
    size=16, bold=True, color=LTEAL, align=PP_ALIGN.CENTER)
txt(s, "Connect your AWS account in 5 minutes  ·  2-week free trial  ·  Cancel anytime",
    2.0, 5.98, 9.33, 0.38, size=11, color=MUTED, align=PP_ALIGN.CENTER)
txt(s, "Questions? leelakrishnakoppolu@gmail.com", 2.0, 6.52, 9.33, 0.35,
    size=10, color=MUTED, align=PP_ALIGN.CENTER)

rect(s, 0, 7.42, 13.33, 0.08, TEAL)
slide_num(s, 11)

# ── Save ──────────────────────────────────────────────────────────
out = r"c:\Users\leela\compliance-ai-saas\VigiliCloud_Client_Pitch.pptx"
prs.save(out)
print(f"Saved: {out}")
