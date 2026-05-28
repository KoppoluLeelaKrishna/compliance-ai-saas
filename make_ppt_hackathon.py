# VigiliCloud - Hackathon & Fortune 500 Pitch Deck
# Run: C:\Users\leela\anaconda3\python.exe make_ppt_hackathon.py

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colour palette ───────────────────────────────────────────────────
BLACK   = RGBColor(0x06, 0x06, 0x06)
DARK    = RGBColor(0x0d, 0x0d, 0x0d)
CARD    = RGBColor(0x16, 0x16, 0x16)
CARD2   = RGBColor(0x1a, 0x1a, 0x1a)
GREEN   = RGBColor(0x10, 0xb9, 0x81)   # emerald-500
LGREEN  = RGBColor(0x34, 0xd3, 0x99)   # emerald-400
DKGREEN = RGBColor(0x06, 0x2a, 0x1e)
WHITE   = RGBColor(0xff, 0xff, 0xff)
GREY    = RGBColor(0xa3, 0xa3, 0xa3)
DKGREY  = RGBColor(0x52, 0x52, 0x52)
RED     = RGBColor(0xef, 0x44, 0x44)
ORANGE  = RGBColor(0xf9, 0x73, 0x16)
YELLOW  = RGBColor(0xea, 0xb3, 0x08)
BLUE    = RGBColor(0x38, 0xbc, 0xf8)
PURPLE  = RGBColor(0xa7, 0x8b, 0xfa)
GOLD    = RGBColor(0xfb, 0xbf, 0x24)

TOTAL_SLIDES = 15

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# ── Helpers ──────────────────────────────────────────────────────────
def bg(slide, color=DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def bordered_box(slide, l, t, w, h, fill_color, border_color, bw=0.8):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill_color
    s.line.color.rgb = border_color; s.line.width = Pt(bw)
    return s

def txt(slide, text, l, t, w, h, size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return tb

def accent_bar(slide, t=0.0, color=GREEN):
    b = box(slide, 0, t, 13.33, 0.06, color)
    return b

def slide_num(slide, n):
    txt(slide, f"{n} / {TOTAL_SLIDES}", 12.3, 7.12, 0.9, 0.3, size=9, color=DKGREY, align=PP_ALIGN.RIGHT)

def tag(slide, label, l, t, color=GREEN, text_color=BLACK):
    w = len(label) * 0.092 + 0.3
    box(slide, l, t, w, 0.28, color)
    txt(slide, label, l + 0.1, t + 0.03, w - 0.12, 0.24, size=9, bold=True, color=text_color, align=PP_ALIGN.CENTER)

def section_header(slide, title, subtitle=None):
    box(slide, 0, 0, 13.33, 1.22, RGBColor(0x06, 0x2a, 0x1e))
    accent_bar(slide, 1.22, GREEN)
    txt(slide, title, 0.45, 0.08, 11.5, 0.62, size=22, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, 0.45, 0.72, 11.5, 0.42, size=11, color=GREY)

def card(slide, l, t, w, h, fill=CARD, border=GREEN, bw=0.6):
    return bordered_box(slide, l, t, w, h, fill, border, bw)


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE / HERO
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, BLACK)
accent_bar(s, 0.0, GREEN)

# Left hero text
txt(s, "V", 0.45, 0.55, 0.72, 0.72, size=30, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
bordered_box(s, 0.45, 0.55, 0.72, 0.72, DKGREEN, GREEN, 1.0)
txt(s, "V", 0.45, 0.55, 0.72, 0.72, size=30, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

txt(s, "VigiliCloud", 1.3, 0.58, 6, 0.62, size=32, bold=True, color=WHITE)
txt(s, "AWS Security Intelligence Platform", 1.3, 1.25, 7, 0.42, size=14, color=GREY)

box(s, 1.3, 1.78, 5.5, 0.04, GREEN)

txt(s, "Detect Misconfigurations.  Get AI-Powered Fix Guidance.  Prove Compliance.",
    1.3, 1.92, 8.5, 0.52, size=14, bold=True, color=LGREEN)

bullets = [
    "  10 critical AWS security checks — scans complete in under 2 minutes",
    "  Claude AI executive summaries with prioritised remediation steps",
    "  SOC2 / ISO 27001 evidence exports — one click",
    "  $0 stored credentials — read-only IAM role via AWS STS",
]
for i, b in enumerate(bullets):
    txt(s, b, 1.3, 2.58 + i * 0.44, 7.5, 0.4, size=11.5, color=GREY)

txt(s, "Built solo. Deployed to production. Generating revenue.",
    1.3, 4.52, 7.5, 0.38, size=12, bold=True, color=LGREEN)

# Presenter line
txt(s, "Leela Krishna Koppolu  ·  vigilicloud.com  ·  leelakrishnakoppolu@gmail.com  ·  2026",
    0.45, 6.9, 9, 0.38, size=10, color=DKGREY)

# Right panel — live scan preview
rp = bordered_box(s, 9.5, 0.45, 3.6, 6.6, CARD, GREEN, 0.8)
txt(s, "LIVE SCAN RESULTS", 9.65, 0.58, 3.3, 0.32, size=9, bold=True, color=GREEN)
box(s, 9.65, 0.92, 3.3, 0.03, DKGREY)
findings = [
    ("🪣  S3 Bucket Public",     "CRITICAL", RED),
    ("🔑  Root Access Keys",     "CRITICAL", RED),
    ("🔐  IAM Over-Permission",  "HIGH",     ORANGE),
    ("📱  MFA Not Enforced",     "HIGH",     ORANGE),
    ("🛡️   Open Security Groups","HIGH",     ORANGE),
    ("🗄️   RDS Unencrypted",    "HIGH",     ORANGE),
    ("💾  EBS Encryption Gap",   "MEDIUM",  YELLOW),
    ("📋  CloudTrail Inactive",  "MEDIUM",  YELLOW),
    ("🌐  VPC Flow Logs Off",    "MEDIUM",  YELLOW),
    ("🔄  KMS No Auto-Rotate",   "MEDIUM",  YELLOW),
]
for i, (name, sev, col) in enumerate(findings):
    y = 1.05 + i * 0.5
    txt(s, name, 9.65, y, 2.3, 0.42, size=9.5, color=WHITE)
    txt(s, sev, 11.95, y + 0.06, 1.1, 0.3, size=8.5, bold=True, color=col, align=PP_ALIGN.RIGHT)

txt(s, "✓  Scan Time: 1m 42s   ·   10/10 Checks", 9.65, 6.05, 3.3, 0.3, size=9, color=LGREEN)

slide_num(s, 1)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — THE HOOK
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, BLACK)
accent_bar(s, 0.0)

txt(s, "Right now, as you read this slide...",
    0.8, 0.6, 11.73, 0.55, size=20, color=GREY, align=PP_ALIGN.CENTER, italic=True)

txt(s, "78% of AWS accounts have at least\none CRITICAL security misconfiguration.",
    0.5, 1.3, 12.33, 1.4, size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(s, "Is yours one of them?",
    0.5, 2.85, 12.33, 0.65, size=28, bold=True, color=LGREEN, align=PP_ALIGN.CENTER)

box(s, 3.5, 3.65, 6.33, 0.05, GREEN)

stats = [
    ("$4.5M",       "Average cost of\na cloud data breach\n(IBM 2024)",         RED),
    ("2 days",      "Average time to\ndetect a CRITICAL\nmisconfiguration",      ORANGE),
    ("82%",         "Of breaches involve\nmisconfigured cloud\nresources",        YELLOW),
    ("0",           "Seconds it takes an\nattacker to find\nan open S3 bucket",  LGREEN),
]
for i, (val, lbl, col) in enumerate(stats):
    x = 0.55 + i * 3.1
    bordered_box(s, x, 3.88, 2.85, 2.6, CARD, col, 1.0)
    txt(s, val, x + 0.15, 4.05, 2.55, 0.72, size=34, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(s, lbl, x + 0.15, 4.82, 2.55, 0.95, size=11, color=GREY, align=PP_ALIGN.CENTER)

txt(s, "VigiliCloud finds every one of these issues in under 2 minutes — and tells you exactly how to fix them.",
    0.5, 6.72, 12.33, 0.45, size=12, bold=True, color=LGREEN, align=PP_ALIGN.CENTER)

slide_num(s, 2)

# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — THE PROBLEM (4 PAIN POINTS)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK)
section_header(s, "The Problem", "AWS misconfigurations are silently destroying businesses — today")

pains = [
    ("No Visibility",      RED,    "🚨",
     "Teams don't know what's exposed.\nNo central view of who has admin access,\nwhat's publicly accessible, or what's unencrypted.",
     "68% of companies don't know their full AWS security posture"),
    ("No Remediation Path",ORANGE, "🔧",
     "Existing tools give you alerts but\nno guidance. Engineers know something is\nwrong but don't know exactly how to fix it.",
     "Average remediation time: 14 days when guidance is missing"),
    ("Manual Audits Fail", YELLOW, "⏱️",
     "Manual security reviews take 2–3 days\nand still miss things. They're expensive,\nerror-prone, and can't scale.",
     "A DevOps engineer costs $150/hr — 2-day audit = $2,400"),
    ("Compliance Evidence", BLUE,  "📋",
     "SOC2 and ISO 27001 audits require\nevidence packages that take weeks to\ncollect manually — costing enterprise deals.",
     "Companies lose deals because they can't produce compliance evidence fast enough"),
]
for i, (title, col, icon, desc, stat) in enumerate(pains):
    ci = i % 2; row = i // 2
    x = 0.45 + ci * 6.44
    y = 1.42 + row * 2.9
    bordered_box(s, x, y, 6.14, 2.72, CARD, col, 1.0)
    box(s, x, y, 6.14, 0.06, col)
    txt(s, icon, x + 0.2, y + 0.18, 0.55, 0.55, size=22)
    txt(s, title, x + 0.85, y + 0.22, 5.1, 0.42, size=15, bold=True, color=WHITE)
    txt(s, desc, x + 0.2, y + 0.82, 5.74, 0.95, size=10.5, color=GREY)
    box(s, x + 0.2, y + 1.9, 5.74, 0.55, RGBColor(0x12, 0x12, 0x12))
    txt(s, f"📊  {stat}", x + 0.32, y + 1.98, 5.5, 0.38, size=9.5, color=col, italic=True)

slide_num(s, 3)

# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — WHAT I BUILT: THE SOLUTION
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK)
section_header(s, "What I Built — VigiliCloud", "A full production SaaS: from idea to live product to paying customers")

txt(s, "Connect your AWS account → Run a full security scan → Get AI-powered fix guidance → Export compliance evidence.",
    0.45, 1.42, 12.43, 0.42, size=12.5, color=GREY)

steps = [
    ("01", "Connect AWS\nAccount", "Create a read-only\nIAM role. Paste ARN.\n5 minutes, once.",     GREEN, "🔗", "No stored\ncredentials"),
    ("02", "Run Security\nScan",   "10 CIS checks run\nin parallel across\nyour entire account.",    ORANGE,"🔍", "~2 min\ncomplete"),
    ("03", "AI Analysis\n(Claude)","Executive summary,\nrisk priority, plain\nEnglish explanation.",  BLUE,  "🤖", "Powered by\nClaude AI"),
    ("04", "Fix What's\nWrong",    "Exact AWS Console\npath, CLI commands,\nstep-by-step guide.",    LGREEN,"🔧", "Fix time\n<1 day avg"),
    ("05", "Export\nEvidence",     "CSV / JSON exports\nfor SOC2, ISO 27001\naudit packages.",       PURPLE,"📋", "One-click\nexport"),
]
for i, (num, title, desc, col, icon, badge) in enumerate(steps):
    x = 0.42 + i * 2.52
    bordered_box(s, x, 2.05, 2.38, 4.55, CARD, col, 0.8)
    box(s, x, 2.05, 2.38, 0.06, col)
    txt(s, num, x + 0.15, 2.15, 0.55, 0.5, size=22, bold=True, color=col)
    txt(s, icon, x + 0.85, 2.18, 0.55, 0.44, size=20)
    txt(s, title, x + 0.18, 2.72, 2.05, 0.52, size=12.5, bold=True, color=WHITE)
    txt(s, desc, x + 0.18, 3.3, 2.05, 0.95, size=10, color=GREY)
    box(s, x + 0.15, 4.98, 2.1, 0.5, DKGREEN)
    txt(s, badge, x + 0.15, 4.98, 2.1, 0.5, size=9.5, bold=True, color=LGREEN, align=PP_ALIGN.CENTER)

    if i < 4:
        txt(s, "→", x + 2.42, 3.95, 0.26, 0.38, size=16, bold=True, color=DKGREY, align=PP_ALIGN.CENTER)

txt(s, "✦  Live at app.vigilicloud.com  ·  Setup in 5 minutes  ·  No credit card for trial  ·  Works with any AWS account",
    0.45, 6.82, 12.43, 0.38, size=11, color=DKGREY, align=PP_ALIGN.CENTER)

slide_num(s, 4)

# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — TECHNICAL ARCHITECTURE (WHAT I BUILT)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK)
section_header(s, "Technical Architecture — Built Solo in Production", "Full-stack SaaS: FastAPI + Next.js + Claude AI + AWS boto3")

# Architecture flow
layers = [
    ("Next.js 16 UI", "React, Tailwind 4\nDark theme, emerald\naccents. Deployed\non Render CDN.", BLUE, 0.45, 1.55, 3.0, 4.8),
    ("FastAPI Backend", "Python REST API\nSession auth, rate\nlimiting, billing,\nscan orchestration.", LGREEN, 3.75, 1.55, 3.0, 4.8),
    ("Compliance Worker", "10 boto3 checks\nrun in parallel:\nS3, IAM, EC2, RDS,\nVPC, KMS checks.", ORANGE, 7.05, 1.55, 3.0, 4.8),
    ("Claude AI (Haiku)", "Anthropic SDK\nExecutive summary\nper scan. Priority\nrisk analysis.", PURPLE, 10.1, 1.55, 2.95, 4.8),
]
for title, desc, col, x, y, w, h in layers:
    bordered_box(s, x, y, w, h, CARD, col, 1.0)
    box(s, x, y, w, 0.06, col)
    txt(s, title, x + 0.18, y + 0.15, w - 0.3, 0.42, size=13, bold=True, color=col)
    txt(s, desc, x + 0.18, y + 0.65, w - 0.3, 1.3, size=10.5, color=GREY)

    # Arrow
    if x < 10:
        txt(s, "→", x + w + 0.08, y + 2.1, 0.45, 0.5, size=18, bold=True, color=DKGREY, align=PP_ALIGN.CENTER)

# Infrastructure row
box(s, 0.45, 6.52, 12.6, 0.05, DKGREY)
infra = [
    ("PostgreSQL\n(Render)", GREY, 0.45),
    ("Razorpay\nBilling", GREY, 3.0),
    ("Resend Email\nAlerts", GREY, 5.7),
    ("AWS STS\nAssumeRole", GREY, 8.35),
    ("GitHub → Render\nAuto-deploy", GREY, 10.85),
]
txt(s, "Infrastructure:", 0.45, 6.6, 2.0, 0.35, size=9.5, color=DKGREY, bold=True)
for label, col, x in infra:
    txt(s, label, x + 0.05, 6.55, 2.3, 0.4, size=9, color=DKGREY, align=PP_ALIGN.CENTER)

# Built facts
facts_row = [
    ("FastAPI", "Python backend"),
    ("Next.js 16", "React frontend"),
    ("Claude Haiku", "AI analysis"),
    ("10 checks", "CIS Benchmark"),
    ("Postgres", "Production DB"),
    ("Solo-built", "End-to-end"),
]
y_f = 7.0
for i, (val, lbl) in enumerate(facts_row):
    x = 0.45 + i * 2.15
    txt(s, val, x, y_f, 2.0, 0.28, size=10, bold=True, color=LGREEN)
    txt(s, lbl, x, y_f + 0.26, 2.0, 0.22, size=8.5, color=DKGREY)

slide_num(s, 5)

# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — 10 SECURITY CHECKS
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK)
section_header(s, "10 Critical Security Checks — CIS AWS Benchmark", "Every check maps to SOC2, ISO 27001, and PCI DSS compliance requirements")

checks = [
    ("🪣", "S3 Public Access",   "Public buckets, ACLs,\nexposed policies",         "CRITICAL", RED),
    ("🔑", "Root Access Keys",   "Active root keys —\nmost critical risk",           "CRITICAL", RED),
    ("🔐", "IAM Permissions",    "Over-privileged roles,\nunnecessary admin access",  "HIGH",     ORANGE),
    ("📱", "MFA Enforcement",    "Missing MFA on users\nand root account",            "HIGH",     ORANGE),
    ("🛡️", "Security Groups",   "EC2 ports open to\n0.0.0.0/0 internet",            "HIGH",     ORANGE),
    ("🗄️", "RDS Encryption",    "Unencrypted DBs,\npublicly accessible RDS",        "HIGH",     ORANGE),
    ("💾", "EBS Encryption",     "Unencrypted volumes,\nmissing default encryption",  "MEDIUM",   YELLOW),
    ("📋", "CloudTrail Logging", "Inactive trails,\nno log validation",              "MEDIUM",   YELLOW),
    ("🌐", "VPC Flow Logs",      "Network traffic\nnot logged",                      "MEDIUM",   YELLOW),
    ("🔄", "KMS Key Rotation",   "Customer keys not\nauto-rotating",                 "MEDIUM",   YELLOW),
]
cols = 5
for i, (icon, name, desc, sev, col) in enumerate(checks):
    row = i // cols; ci = i % cols
    x = 0.35 + ci * 2.55
    y = 1.55 + row * 2.72
    bordered_box(s, x, y, 2.4, 2.55, CARD, col, 0.6)
    txt(s, icon, x + 0.15, y + 0.18, 0.5, 0.48, size=18)
    txt(s, sev, x + 0.72, y + 0.24, 1.52, 0.28, size=8, bold=True, color=col)
    txt(s, name, x + 0.15, y + 0.68, 2.1, 0.42, size=11, bold=True, color=WHITE)
    txt(s, desc, x + 0.15, y + 1.15, 2.1, 0.75, size=9.5, color=GREY)

slide_num(s, 6)

# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — AI-POWERED ANALYSIS (CLAUDE)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK)
section_header(s, "AI-Powered Security Analysis — Powered by Claude", "Every scan generates an AI executive summary — plain English, prioritised, actionable")

# Left: before AI
bordered_box(s, 0.45, 1.52, 5.9, 5.5, CARD, RED, 0.8)
box(s, 0.45, 1.52, 5.9, 0.06, RED)
txt(s, "❌  WITHOUT AI (Traditional Tools)", 0.65, 1.62, 5.5, 0.38, size=12, bold=True, color=RED)

before_items = [
    ("Finding #1", "arn:aws:s3:::my-company-data — PublicAccessBlockConfiguration.BlockPublicAcls: false"),
    ("Finding #2", "arn:aws:iam::123456789:user/john — MFAActive: false, PasswordLastUsed: 2024-01-05"),
    ("Finding #3", "arn:aws:ec2::sg-0abc123 — IpRanges: [{'CidrIp': '0.0.0.0/0'}] port 22"),
    ("Finding #4", "arn:aws:rds::db-prod — PubliclyAccessible: true, StorageEncrypted: false"),
    ("Finding #5", "arn:aws:iam::root — AccessKeys[0].Status: Active, CreateDate: 2021-06-15"),
    ("...", "7 more raw API responses..."),
]
for i, (label, text) in enumerate(before_items):
    txt(s, label, 0.65, 2.1 + i * 0.68, 1.2, 0.3, size=8.5, bold=True, color=RED)
    txt(s, text, 0.65, 2.36 + i * 0.68, 5.55, 0.28, size=8, color=DKGREY, italic=True)

# Right: with AI (Claude)
bordered_box(s, 6.55, 1.52, 6.45, 5.5, DKGREEN, LGREEN, 1.2)
box(s, 6.55, 1.52, 6.45, 0.06, LGREEN)
txt(s, "✅  WITH CLAUDE AI (VigiliCloud)", 6.75, 1.62, 6.05, 0.38, size=12, bold=True, color=LGREEN)

ai_summary = (
    "EXECUTIVE SUMMARY — Critical Action Required\n\n"
    "Your AWS account has 5 CRITICAL and 4 HIGH severity\n"
    "findings that expose you to immediate data breach risk.\n\n"
    "TOP PRIORITY — FIX IMMEDIATELY:\n"
    "1. Root access keys are ACTIVE — these must be deleted\n"
    "   today. Active root keys give anyone who finds them\n"
    "   total control of your AWS account.\n\n"
    "2. S3 bucket 'my-company-data' is PUBLIC — all files\n"
    "   are accessible to anyone on the internet right now.\n\n"
    "3. RDS 'db-prod' is publicly accessible and unencrypted.\n"
    "   Customer data may be exposed.\n\n"
    "ESTIMATED RISK: $2.1M exposure based on data volume.\n"
    "Estimated fix time with guidance: 4 hours."
)
txt(s, ai_summary, 6.75, 2.08, 6.05, 4.4, size=10, color=WHITE)

# Arrow in middle
txt(s, "→\nClaude\nAI", 6.0, 3.5, 0.5, 0.9, size=10, bold=True, color=LGREEN, align=PP_ALIGN.CENTER)

txt(s, "Claude AI turns raw API data into an executive brief any CTO or auditor can act on immediately.",
    0.45, 7.12, 12.43, 0.28, size=10, color=GREY, align=PP_ALIGN.CENTER, italic=True)

slide_num(s, 7)

# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — KEY FEATURES BUILT
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK)
section_header(s, "Key Features — All Live in Production", "Everything your security and compliance team needs — nothing extra")

features = [
    ("🤖", "Claude AI Analysis",   "Generates plain-English executive\nsummary per scan. Priority scoring.\nRisk explanation for non-technical stakeholders.", LGREEN),
    ("📧", "Instant Email Alerts", "CRITICAL finding? You're notified\nby email within seconds of scan\ncompletion. Zero missed incidents.",                  LGREEN),
    ("📁", "Compliance Exports",   "One-click CSV or JSON export of\nall findings with resource IDs,\nseverity, and fix steps. SOC2-ready.",                  LGREEN),
    ("⏱️", "2-Minute Full Scans", "All 10 CIS checks run in parallel\nvia boto3. Schedule daily auto-scans\nfor continuous posture monitoring.",             LGREEN),
    ("🔒", "Zero Stored Creds",   "Read-only IAM role via AWS STS.\nCredentials never stored — they\nexpire after each scan automatically.",                  LGREEN),
    ("🗺️", "Guided Fix Steps",    "Every finding includes: exact AWS\nConsole path, CLI command, and\nstep-by-step fix instructions.",                        LGREEN),
]
for i, (icon, title, desc, col) in enumerate(features):
    ci = i % 3; row = i // 2
    x = 0.45 + ci * 4.25
    y = 1.52 + (i // 3) * 2.72
    bordered_box(s, x, y, 3.95, 2.52, CARD, GREEN, 0.5)
    txt(s, icon, x + 0.2, y + 0.22, 0.6, 0.52, size=22)
    txt(s, title, x + 0.2, y + 0.82, 3.55, 0.42, size=13, bold=True, color=WHITE)
    txt(s, desc, x + 0.2, y + 1.3, 3.55, 0.98, size=10.5, color=GREY)

slide_num(s, 8)

# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — MARKET OPPORTUNITY
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK)
section_header(s, "Market Opportunity", "A $50B market with a massive SMB gap — and perfect timing")

markets = [
    ("TAM", "$50B+",  "Total cloud security\nmarket (2026, +15% YoY)",         PURPLE, 0.45),
    ("SAM", "$8B",    "SMB & mid-market AWS\nsecurity — our segment",           BLUE,   4.75),
    ("SOM", "$50M",   "Realistic 5-year revenue\ntarget with current GTM",      LGREEN, 9.05),
]
for label, val, desc, col, x in markets:
    bordered_box(s, x, 1.55, 4.05, 3.35, CARD, col, 1.2)
    box(s, x + 0.2, 1.75, 0.9, 0.45, col)
    txt(s, label, x + 0.2, 1.75, 0.9, 0.45, size=13, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    txt(s, val, x + 0.18, 2.28, 3.7, 0.92, size=40, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(s, desc, x + 0.18, 3.28, 3.7, 0.5, size=11, color=GREY, align=PP_ALIGN.CENTER)

# Why now
bordered_box(s, 0.45, 5.1, 12.43, 2.1, CARD, GOLD, 1.0)
txt(s, "⚡  Why NOW Is the Perfect Moment", 0.65, 5.22, 12.0, 0.38, size=13, bold=True, color=GOLD)
timing = [
    ("AWS adoption explosion", "60% of SMBs moved to cloud post-2020 — millions of new accounts with no security team"),
    ("SOC2 now mandatory",     "Enterprise deals now REQUIRE SOC2 evidence — compliance demand is at all-time high"),
    ("AI makes this possible", "Claude AI enables sophisticated security analysis accessible to non-security teams"),
    ("Zero SMB competitors",   "No AWS scanner exists at $99–$999/mo with AI guidance — the gap is wide open"),
]
for i, (title, desc) in enumerate(timing):
    ci = i % 2; row = i // 2
    x = 0.65 + ci * 6.3
    y = 5.7 + row * 0.55
    txt(s, f"✦ {title}:", x, y, 2.0, 0.3, size=10, bold=True, color=GOLD)
    txt(s, desc, x + 2.1, y, 4.05, 0.3, size=10, color=GREY)

slide_num(s, 9)

# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — REVENUE MODEL
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK)
section_header(s, "Revenue Model — SaaS Subscription", "Predictable recurring revenue across 3 tiers — 2-week free trial on all plans")

plans = [
    ("STARTER",   "$99/mo",  "₹8,299/mo",  "Solo consultants\n& small teams",
     ["Up to 3 AWS accounts", "All 10 security checks", "CSV / JSON exports",
      "Email alerts on CRITICAL", "Fix guidance for every finding"],
     GREY, False),
    ("PRO  ★",   "$299/mo", "₹24,999/mo", "Teams managing\nmultiple environments",
     ["Up to 10 AWS accounts", "Everything in Starter", "Claude AI analysis per scan",
      "Scheduled daily auto-scans", "Priority email support"],
     LGREEN, True),
    ("MSP",      "$999/mo", "₹83,499/mo", "Agencies & managed\nservice providers",
     ["Unlimited AWS accounts", "Everything in Pro", "Multi-customer workflows",
      "White-label ready reports", "High-volume scanning"],
     PURPLE, False),
]
for i, (name, usd, inr, desc, feats, col, hot) in enumerate(plans):
    x = 0.45 + i * 4.25
    fill = DKGREEN if hot else CARD
    border = LGREEN if hot else col
    bordered_box(s, x, 1.52, 3.95, 5.3, fill, border, 1.2 if hot else 0.7)
    if hot:
        box(s, x + 0.9, 1.34, 2.2, 0.28, LGREEN)
        txt(s, "★  MOST POPULAR", x + 0.9, 1.36, 2.2, 0.24, size=8.5, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    txt(s, name, x + 0.2, 1.65, 3.55, 0.42, size=14, bold=True, color=col if not hot else WHITE)
    txt(s, usd, x + 0.2, 2.1, 3.55, 0.58, size=28, bold=True, color=LGREEN if hot else WHITE)
    txt(s, inr, x + 0.2, 2.7, 3.55, 0.3, size=10, color=DKGREY)
    txt(s, desc, x + 0.2, 3.05, 3.55, 0.42, size=10, color=GREY)
    for j, f in enumerate(feats):
        txt(s, f"✓  {f}", x + 0.2, 3.58 + j * 0.44, 3.55, 0.4, size=10,
            color=LGREEN if hot else GREY)

# ROI callout
bordered_box(s, 0.45, 7.0, 12.43, 0.38, DKGREEN, GREEN, 0.8)
txt(s, "💎  ROI: One prevented S3 breach saves avg $4.5M. VigiliCloud Pro = $299/mo. That's 1,256x ROI on a single prevented incident.",
    0.65, 7.06, 12.0, 0.28, size=10.5, bold=True, color=LGREEN)

slide_num(s, 10)

# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — REVENUE PROJECTIONS
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK)
section_header(s, "Revenue Projections", "Conservative model based on direct outreach + ProductHunt launch + content marketing")

monthly = [
    ("M1",  0,   0,      "Pre-launch"),
    ("M2",  3,   450,    "Early access"),
    ("M3",  8,   1200,   "Launch"),
    ("M4",  15,  2250,   "Grow"),
    ("M5",  25,  3750,   "Grow"),
    ("M6",  38,  5700,   "Scale"),
    ("M7",  55,  8250,   "Scale"),
    ("M8",  72,  10800,  "Scale"),
    ("M9",  90,  13500,  "Optimize"),
    ("M10", 108, 16200,  "Optimize"),
    ("M11", 125, 18750,  "Steady"),
    ("M12", 145, 21750,  "End Y1"),
]

max_mrr = 21750
chart_h = 3.7
chart_y = 1.55
chart_x = 0.45
bar_area_w = 9.5
bar_w = bar_area_w / len(monthly) - 0.05

bordered_box(s, chart_x, chart_y, bar_area_w, chart_h, CARD, DKGREY, 0.5)

for i, (month, customers, mrr, note) in enumerate(monthly):
    bx = chart_x + i * (bar_w + 0.05) + 0.08
    if mrr > 0:
        bh = (mrr / max_mrr) * (chart_h - 0.5)
        by = chart_y + chart_h - bh - 0.28
        if mrr > 10000:
            col = GREEN
        elif mrr > 3000:
            col = LGREEN
        else:
            col = RGBColor(0x86, 0xef, 0xac)
        bar = box(s, bx, by, bar_w, bh, col)
        if mrr >= 5700:
            label = f"${mrr//1000}K"
            txt(s, label, bx, by - 0.28, bar_w, 0.25, size=7.5, bold=True, color=LGREEN, align=PP_ALIGN.CENTER)
    txt(s, month, bx, chart_y + chart_h - 0.26, bar_w, 0.24, size=8, color=DKGREY, align=PP_ALIGN.CENTER)

txt(s, "Monthly Recurring Revenue (USD)", chart_x + 2, 5.42, 6, 0.28, size=9, color=DKGREY, align=PP_ALIGN.CENTER)

# Right summary
bordered_box(s, 10.15, 1.55, 2.75, 4.78, CARD, GREEN, 0.8)
txt(s, "Year 1 Summary", 10.3, 1.68, 2.45, 0.35, size=11, bold=True, color=GREEN)
y1 = [
    ("End-MRR",    "$21.7K"),
    ("End-ARR",    "$261K"),
    ("Customers",  "145"),
    ("Avg ARPU",   "$150/mo"),
    ("Trial→Paid", "35%"),
    ("Churn",      "<5%"),
]
for i, (lbl, val) in enumerate(y1):
    ys = 2.12 + i * 0.62
    txt(s, lbl, 10.3, ys, 1.2, 0.3, size=9.5, color=GREY)
    txt(s, val, 11.52, ys, 1.25, 0.3, size=10.5, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
    box(s, 10.3, ys + 0.42, 2.45, 0.02, DKGREY)

# 3-year outlook
txt(s, "3-Year Outlook", 0.45, 5.65, 9.5, 0.35, size=11.5, bold=True, color=WHITE)
years = [
    ("Year 1", "145 customers",   "$261K ARR",  GREY),
    ("Year 2", "500 customers",   "$1.2M ARR",  BLUE),
    ("Year 3", "1,500 customers", "$5.4M ARR",  LGREEN),
]
for i, (yr, cust, arr, col) in enumerate(years):
    x3 = 0.45 + i * 3.15
    bordered_box(s, x3, 6.08, 3.0, 1.05, CARD, col, 0.8)
    box(s, x3, 6.08, 3.0, 0.05, col)
    txt(s, yr,   x3 + 0.15, 6.17, 1.0, 0.3,  size=10, bold=True, color=col)
    txt(s, cust, x3 + 0.15, 6.47, 1.5, 0.28, size=9.5, color=GREY)
    txt(s, arr,  x3 + 1.72, 6.15, 1.1, 0.3,  size=11, bold=True, color=col, align=PP_ALIGN.RIGHT)

slide_num(s, 11)

# ═══════════════════════════════════════════════════════════════
# SLIDE 12 — COMPETITIVE LANDSCAPE
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK)
section_header(s, "Competitive Landscape", "VigiliCloud is the only tool combining simplicity + AI guidance + compliance exports at SMB pricing")

headers = ["Feature",         "VigiliCloud",   "AWS Security Hub",    "Prowler (OSS)",   "Lacework / Wiz",   "Manual Audit"]
col_w   = [2.55, 1.75, 1.75, 1.75, 1.75, 1.55]
col_x   = [0.28]
for w in col_w[:-1]: col_x.append(col_x[-1] + w)

for i, (h, cx, cw) in enumerate(zip(headers, col_x, col_w)):
    bg_col = GREEN if i == 1 else RGBColor(0x1a, 0x1a, 0x1a)
    b = box(s, cx, 1.55, cw - 0.06, 0.52, bg_col)
    tc = BLACK if i == 1 else GREY
    txt(s, h, cx + 0.08, 1.6, cw - 0.18, 0.44, size=9, bold=True, color=tc, align=PP_ALIGN.CENTER)

rows = [
    ("Setup time",         "5 minutes",    "Hours + config",   "CLI expertise",  "Days + POC",       "2–3 days"),
    ("Price / month",      "$99–$999",     "$30–$1,000+",      "Free (DIY)",     "$5K–$50K+",        "Your time"),
    ("Fix guidance",       "✔  Full",      "✘  None",          "✘  None",        "△  Limited",       "✘  Manual"),
    ("AI analysis",        "✔  Claude AI", "✘  None",          "✘  None",        "△  Some",          "✘  None"),
    ("No expertise req.",  "✔  Anyone",    "✘  AWS expert",    "✘  Dev only",    "✘  Security eng.", "✘  Expert"),
    ("Compliance export",  "✔  CSV/JSON",  "△  Partial",       "✘  Limited",     "✔  Enterprise",    "✔  Manual"),
    ("SMB-accessible",     "✔  YES",       "△  Complex",       "✘  Dev tool",    "✘  Enterprise",    "✔  Yes"),
]
tick_map = {"✔": LGREEN, "✘": RED, "△": YELLOW}
for ri, row in enumerate(rows):
    y = 2.14 + ri * 0.58
    row_bg = CARD if ri % 2 == 0 else CARD2
    for ci, (cell, cx, cw) in enumerate(zip(row, col_x, col_w)):
        b = box(s, cx, y, cw - 0.06, 0.52, row_bg)
        if ci == 1: b.fill.fore_color.rgb = DKGREEN
        first = cell[:1]
        tc = tick_map.get(first, WHITE if ci == 0 else GREY)
        txt(s, cell, cx + 0.08, y + 0.1, cw - 0.16, 0.36, size=9.5,
            bold=(ci == 1), color=tc, align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)

bordered_box(s, 0.28, 6.34, 12.77, 0.42, DKGREEN, LGREEN, 0.8)
txt(s, "★  VigiliCloud is the ONLY tool with AI analysis + full fix guidance + compliance exports at under $1,000/month",
    0.48, 6.42, 12.4, 0.3, size=11, bold=True, color=LGREEN)

slide_num(s, 12)

# ═══════════════════════════════════════════════════════════════
# SLIDE 13 — FUTURE ROADMAP
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK)
section_header(s, "Future Roadmap — From AWS Scanner to Cloud Security Platform", "4 phases to become the #1 cloud security platform for SMBs globally")

phases = [
    ("Phase 1\nTODAY",        "AWS Security Scanner",
     ["10 CIS checks (live)", "Claude AI analysis (live)", "Compliance exports (live)", "Billing + subscriptions"],
     "LIVE", LGREEN, DKGREEN),
    ("Phase 2\nQ3–Q4 2026",   "Multi-Cloud + Alerts",
     ["Azure & GCP support", "Slack / Teams alerts", "Developer API access", "MCP auto-fix proposals"],
     "Q3-Q4 2026", BLUE, RGBColor(0x08, 0x1c, 0x38)),
    ("Phase 3\n2027",         "AI Remediation Engine",
     ["One-click AI fix approval", "Jira / GitHub PR auto-create", "IaC scanning (Terraform/CDK)", "SOC2 Type II evidence"],
     "2027", ORANGE, RGBColor(0x2a, 0x10, 0x00)),
    ("Phase 4\n2027+",        "Enterprise Platform",
     ["White-label MSP mode", "HIPAA / PCI check modules", "SIEM integration", "24/7 continuous monitoring"],
     "VISION", PURPLE, RGBColor(0x16, 0x08, 0x32)),
]
for i, (phase, title, items, status, col, bg_c) in enumerate(phases):
    x = 0.45 + i * 3.22
    bordered_box(s, x, 1.52, 3.05, 5.52, bg_c, col, 1.2)
    box(s, x, 1.52, 3.05, 0.06, col)
    b = box(s, x + 0.18, 1.65, 2.69, 0.3, col)
    txt(s, status.upper(), x + 0.18, 1.66, 2.69, 0.28, size=8.5, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    txt(s, phase, x + 0.18, 2.05, 2.69, 0.5, size=11, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(s, title, x + 0.18, 2.6, 2.69, 0.55, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    box(s, x + 0.35, 3.22, 2.35, 0.04, col)
    for j, item in enumerate(items):
        iy = 3.38 + j * 0.72
        txt(s, f"▶  {item}", x + 0.25, iy, 2.6, 0.62, size=10.5, color=GREY)
    if i < 3:
        txt(s, "→", x + 3.08, 3.8, 0.22, 0.42, size=16, bold=True, color=DKGREY, align=PP_ALIGN.CENTER)

# MCP callout
bordered_box(s, 0.45, 7.12, 12.43, 0.28, RGBColor(0x16, 0x08, 0x32), PURPLE, 0.8)
txt(s, "🚀  Phase 3 MCP Vision: Claude AI proposes fixes → you click APPROVE → Claude executes them automatically → auto-verify confirms resolution",
    0.65, 7.18, 12.0, 0.2, size=10, bold=True, color=PURPLE)

slide_num(s, 13)

# ═══════════════════════════════════════════════════════════════
# SLIDE 14 — WHY THIS MATTERS (FOR FORTUNE 500)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK)
section_header(s, "Why VigiliCloud Matters — For You", "Fortune 500 companies & enterprises: here's where we fit in your ecosystem")

angles = [
    ("🏢", "For Enterprise Procurement",   BLUE,
     "Your portfolio companies, suppliers, and SMB partners run AWS without proper security posture checks. VigiliCloud gives them enterprise-grade security at SMB-accessible pricing — reducing YOUR supply chain risk.",
     "Supply chain security is now a Board-level requirement. ISO 27001 mandates third-party risk management."),
    ("🤝", "For Strategic Partnerships",   LGREEN,
     "White-label VigiliCloud under your brand. Offer AWS security scanning as a value-add service to your SMB clients. Turn our platform into your product — we handle the tech, you own the relationship.",
     "MSP tier supports unlimited accounts. White-label branding ready. Revenue share model available."),
    ("💼", "For Investment / Partnership", GOLD,
     "A solo-built, fully deployed, production SaaS with a clear $50B market, proven tech stack, 3-year revenue model, and a roadmap to AI-powered auto-remediation. Seed-stage with enterprise-grade architecture.",
     "Year 1 target: $261K ARR. Year 3 target: $5.4M ARR. 1,500 customers at $300 avg ARPU."),
    ("🔬", "For Innovation Labs",           ORANGE,
     "VigiliCloud demonstrates what's possible when you combine AWS security expertise with Claude AI. The MCP Phase 2 roadmap — where AI proposes and executes fixes with one-click approval — is a new paradigm in security ops.",
     "Built using Anthropic Claude + AWS STS + FastAPI + Next.js. Full Claude AI integration from day one."),
]
for i, (icon, title, col, desc, stat) in enumerate(angles):
    ci = i % 2; row = i // 2
    x = 0.45 + ci * 6.44
    y = 1.52 + row * 2.95
    bordered_box(s, x, y, 6.14, 2.78, CARD, col, 0.8)
    box(s, x, y, 6.14, 0.06, col)
    txt(s, icon, x + 0.2, y + 0.2, 0.55, 0.52, size=22)
    txt(s, title, x + 0.85, y + 0.24, 5.1, 0.42, size=13, bold=True, color=col)
    txt(s, desc, x + 0.2, y + 0.82, 5.74, 1.05, size=10, color=GREY)
    box(s, x + 0.2, y + 2.0, 5.74, 0.58, CARD2)
    txt(s, f"  {stat}", x + 0.28, y + 2.08, 5.58, 0.42, size=9.5, color=col, italic=True)

slide_num(s, 14)

# ═══════════════════════════════════════════════════════════════
# SLIDE 15 — CALL TO ACTION / CLOSE
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, BLACK)
accent_bar(s, 0.0, GREEN)

txt(s, "The question isn't whether your AWS\naccount has security gaps.",
    0.6, 0.38, 11.9, 1.1, size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(s, "The question is: will you find them before an attacker does?",
    0.6, 1.55, 11.9, 0.6, size=21, bold=True, color=LGREEN, align=PP_ALIGN.CENTER)

box(s, 4.5, 2.28, 4.33, 0.05, GREEN)

# 3 CTA cards
ctas = [
    ("🚀  Try VigiliCloud Free",    "app.vigilicloud.com\n2-week trial — no credit card\nScan your first account in 5 mins", GREEN,   DKGREEN,  "For potential customers"),
    ("🤝  Partner With Us",         "White-label partnership\nMSP & referral programs\nleelakrishnakoppolu@gmail.com",         BLUE,    RGBColor(0x08,0x1a,0x38),  "For agencies & enterprises"),
    ("💬  Let's Talk",              "15-minute call\nDemo, investment, or ideas\ncalendly.com/leelakrishnakoppolu",       GOLD,    RGBColor(0x28,0x1e,0x04),  "For investors & advisors"),
]
for i, (title, detail, col, fill, badge) in enumerate(ctas):
    x = 0.55 + i * 4.15
    bordered_box(s, x, 2.45, 3.9, 2.72, fill, col, 1.5)
    box(s, x, 2.45, 3.9, 0.06, col)
    txt(s, title, x + 0.2, 2.58, 3.5, 0.45, size=14, bold=True, color=col)
    txt(s, detail, x + 0.2, 3.1, 3.5, 0.85, size=11, color=WHITE)
    box(s, x + 0.2, 4.72, 3.5, 0.32, col)
    txt(s, badge, x + 0.2, 4.74, 3.5, 0.28, size=8.5, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

# Stats row
stats_row = [
    ("$50B",  "Market size"),
    ("78%",   "Accounts with CRITICAL issues"),
    ("10",    "Security checks"),
    ("2 min", "Full AWS scan"),
    ("$99",   "Starting price/mo"),
    ("$4.5M", "Avg breach cost prevented"),
]
for i, (val, lbl) in enumerate(stats_row):
    sx = 0.42 + i * 2.16
    txt(s, val, sx, 5.42, 2.0, 0.5, size=20, bold=True, color=LGREEN, align=PP_ALIGN.CENTER)
    txt(s, lbl, sx, 5.92, 2.0, 0.35, size=8.5, color=DKGREY, align=PP_ALIGN.CENTER)

box(s, 0.45, 6.38, 12.43, 0.04, DKGREY)

txt(s, "Leela Krishna Koppolu  ·  leelakrishnakoppolu@gmail.com  ·  app.vigilicloud.com  ·  May 2026",
    0.6, 6.5, 12.13, 0.35, size=11, color=DKGREY, align=PP_ALIGN.CENTER)
txt(s, "VigiliCloud — AWS Security Intelligence for Every Business. Not Just Enterprise.",
    0.6, 6.92, 12.13, 0.38, size=13, bold=True, color=LGREEN, align=PP_ALIGN.CENTER)

slide_num(s, 15)

# ── Save ─────────────────────────────────────────────────────────────
out = r"c:\Users\leela\compliance-ai-saas\VigiliCloud_Hackathon_Pitch.pptx"

prs.save(out)
print(f"Saved: {out}")
