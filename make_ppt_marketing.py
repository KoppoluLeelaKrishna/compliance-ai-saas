# VigiliCloud - Marketing & Growth Deck
# Covers: market opportunity, revenue projections, roadmap, competitor landscape, targets
# Run: C:\Users\leela\anaconda3\python.exe make_ppt_marketing.py
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Color Palette ──────────────────────────────────────────────────
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY    = RGBColor(0xF1, 0xF5, 0xF9)
BORDER   = RGBColor(0xE2, 0xE8, 0xF0)
NAVY     = RGBColor(0x0F, 0x27, 0x6D)
NAVYMD   = RGBColor(0x1D, 0x4E, 0xD8)
TEAL     = RGBColor(0x05, 0x96, 0x69)
LTEAL    = RGBColor(0x10, 0xB9, 0x81)
TEAL_BG  = RGBColor(0xEC, 0xFD, 0xF5)
NAVY_BG  = RGBColor(0xEF, 0xF6, 0xFF)
TEXT     = RGBColor(0x0F, 0x17, 0x2A)
SUBTEXT  = RGBColor(0x47, 0x55, 0x69)
MUTED    = RGBColor(0x94, 0xA3, 0xB8)
RED      = RGBColor(0xDC, 0x26, 0x26)
RED_BG   = RGBColor(0xFE, 0xF2, 0xF2)
ORANGE   = RGBColor(0xEA, 0x58, 0x0C)
ORG_BG   = RGBColor(0xFF, 0xF7, 0xED)
AMBER    = RGBColor(0xD9, 0x77, 0x06)
AMB_BG   = RGBColor(0xFF, 0xFB, 0xEB)
GOLD     = RGBColor(0xF5, 0x9E, 0x0B)
GOLD_BG  = RGBColor(0xFF, 0xFD, 0xEB)
PURPLE   = RGBColor(0x7C, 0x3A, 0xED)
PURP_BG  = RGBColor(0xF5, 0xF3, 0xFF)
SLATE    = RGBColor(0x64, 0x74, 0x8B)
SLT_BG   = RGBColor(0xF8, 0xFA, 0xFC)
DKGREEN  = RGBColor(0x04, 0x6C, 0x4E)

TOTAL = 13

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def bg(slide, color=LGRAY):
    fill = slide.background.fill
    fill.solid(); fill.fore_color.rgb = color


def rect(slide, l, t, w, h, fill_color, border_color=None, border_w=0.75):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill_color
    if border_color:
        s.line.color.rgb = border_color; s.line.width = Pt(border_w)
    else:
        s.line.fill.background()
    return s


def txt(slide, text, l, t, w, h, size=12, bold=False, color=TEXT,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return tb


def hdr(slide, title, subtitle=None):
    rect(slide, 0, 0, 13.33, 1.3, NAVY)
    rect(slide, 0, 1.3, 13.33, 0.045, GOLD)
    txt(slide, title, 0.45, 0.1, 10.5, 0.65, size=23, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, 0.45, 0.77, 11.8, 0.42, size=11, color=RGBColor(0xFF, 0xE8, 0x90))


def slide_num(slide, n):
    txt(slide, f"{n} of {TOTAL}", 12.2, 7.1, 0.9, 0.3, size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def card(slide, l, t, w, h, fill=WHITE, border=BORDER, bw=0.75):
    return rect(slide, l, t, w, h, fill, border, bw)


def divider(slide, t, color=BORDER):
    rect(slide, 0.45, t, 12.43, 0.02, color)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, NAVY)
rect(s, 0, 0, 0.08, 7.5, GOLD)
rect(s, 0, 7.42, 13.33, 0.08, TEAL)

# Logo
logo = rect(s, 0.55, 0.5, 1.0, 1.0, TEAL)
logo.line.fill.background()
txt(s, "V", 0.55, 0.5, 1.0, 1.0, size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(s, "VigiliCloud", 1.75, 0.55, 7, 0.72, size=34, bold=True, color=WHITE)
txt(s, "AWS Security Scanning Platform", 1.75, 1.3, 7, 0.38, size=14, color=LTEAL)

rect(s, 0.55, 1.88, 9.5, 0.05, GOLD)

txt(s, "Growth & Marketing Deck", 0.55, 2.1, 9.5, 0.55, size=20, bold=True, color=RGBColor(0xFF, 0xE8, 0x90))
txt(s, "Market Opportunity  ·  Revenue Projections  ·  Roadmap  ·  Competitor Analysis  ·  Targets",
    0.55, 2.72, 9.5, 0.38, size=12, color=RGBColor(0xBF, 0xDB, 0xFF))

# Key facts panel
facts = [
    ("$50B", "Cloud Security\nMarket (2026)"),
    ("$99/mo", "Starting\nPrice"),
    ("2 min", "Full AWS\nScan Time"),
    ("0", "Stored\nCredentials"),
]
for i, (val, lbl) in enumerate(facts):
    fx = 0.55 + i * 2.42
    card(s, fx, 3.4, 2.2, 1.5, fill=RGBColor(0x1a, 0x3a, 0x7a), border=GOLD, bw=1.0)
    txt(s, val, fx + 0.1, 3.5, 2.0, 0.62, size=24, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    txt(s, lbl, fx + 0.1, 4.12, 2.0, 0.52, size=9.5, color=RGBColor(0xBF, 0xDB, 0xFF), align=PP_ALIGN.CENTER)

txt(s, "Confidential  ·  Leela Krishna Koppolu  ·  leelakrishnakoppolu@gmail.com  ·  vigilicloud.com  ·  May 2026",
    0.55, 6.9, 12.23, 0.38, size=10, color=MUTED, align=PP_ALIGN.CENTER)

slide_num(s, 1)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 2 — MARKET OPPORTUNITY
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
hdr(s, "Market Opportunity", "A massive and growing market — with a massive underserved segment")

markets = [
    ("TAM", "$50B+", "Total cloud security market globally\n(2026, growing 15% YoY)", NAVY, NAVY_BG),
    ("SAM", "$8B",   "SMB & mid-market AWS security —\nour addressable segment",       NAVYMD, NAVY_BG),
    ("SOM", "$50M",  "Realistic 5-year revenue target\nwith current GTM strategy",      TEAL, TEAL_BG),
]
for i, (label, val, desc, col, bg_c) in enumerate(markets):
    x = 0.45 + i * 4.3
    card(s, x, 1.55, 4.05, 3.2, fill=bg_c, border=col, bw=2.0)
    rect(s, x + 0.2, 1.75, 1.0, 0.5, col)
    txt(s, label, x + 0.2, 1.75, 1.0, 0.5, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, val, x + 0.2, 2.35, 3.65, 0.82, size=36, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(s, desc, x + 0.2, 3.22, 3.65, 0.88, size=11, color=SUBTEXT, align=PP_ALIGN.CENTER)

# Why now
card(s, 0.45, 5.0, 12.43, 1.95, fill=WHITE, border=GOLD, bw=1.5)
txt(s, "Why Now — Perfect Market Timing", 0.7, 5.12, 12.0, 0.38, size=13, bold=True, color=NAVY)
timing = [
    "AWS adoption accelerating — 60% of SMBs moved to cloud post-2020",
    "SOC2 / ISO27001 now mandatory for most B2B SaaS deals — compliance pressure is at all-time high",
    "AI-powered tools (like Claude) make sophisticated analysis accessible to non-security teams",
    "No SMB-focused AWS scanner exists at our price point — the market gap is wide open",
]
for i, t in enumerate(timing):
    xi = 0.7 + (i % 2) * 6.25
    yi = 5.58 + (i // 2) * 0.48
    rect(s, xi, yi + 0.1, 0.15, 0.22, GOLD)
    txt(s, t, xi + 0.25, yi, 5.9, 0.38, size=10.5, color=SUBTEXT)

slide_num(s, 2)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 3 — BUSINESS MODEL
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
hdr(s, "Business Model", "Predictable SaaS revenue — subscription-based with clear upsell path")

# Revenue streams
streams = [
    ("Subscription Revenue", "Monthly/annual SaaS subscriptions across 3 tiers. Primary revenue driver. Low churn due to compliance workflow lock-in.",
     NAVY, NAVY_BG, "85% of revenue"),
    ("Upsell — Pro/MSP",     "Consultants and agencies naturally upgrade as they add client accounts. MSP plan has 10x ARPU vs Starter.",
     TEAL, TEAL_BG, "30% upgrade rate"),
    ("Annual Discount",      "20% discount for annual billing. Improves cash flow and reduces monthly churn. High-value customers prefer annual.",
     GOLD, GOLD_BG, "20% discount"),
]
for i, (title, desc, col, bg_c, badge) in enumerate(streams):
    x = 0.45 + i * 4.3
    card(s, x, 1.55, 4.05, 2.5, fill=bg_c, border=col, bw=1.5)
    rect(s, x + 0.2, 1.72, 0.06, 2.15, col)
    txt(s, title, x + 0.4, 1.75, 3.45, 0.38, size=13, bold=True, color=col)
    txt(s, desc, x + 0.4, 2.2, 3.45, 1.12, size=10.5, color=SUBTEXT)
    badge_box = rect(s, x + 0.4, 3.62, 3.45, 0.3, col)
    badge_box.line.fill.background()
    txt(s, badge, x + 0.4, 3.64, 3.45, 0.26, size=9.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Pricing table
txt(s, "Pricing Tiers & Target ARPU", 0.45, 4.28, 12.43, 0.38, size=12.5, bold=True, color=NAVY)
pricing_rows = [
    ("Starter",  "$99/mo",  "₹8,299/mo",  "Consultants, 1–3 accounts", "~40% of customers", MUTED),
    ("Pro",      "$299/mo", "₹24,999/mo", "Teams, 3–10 accounts",       "~45% of customers", NAVYMD),
    ("MSP",      "$999/mo", "₹83,499/mo", "Agencies, 10+ accounts",     "~15% of customers", TEAL),
]
col_xs = [0.45, 2.3, 3.95, 5.65, 8.55, 11.05]
col_ws = [1.8, 1.6, 1.65, 2.85, 2.45, 2.2]
ph = ["Plan", "USD/mo", "INR/mo", "Target segment", "Mix", "Col"]
for ri, (plan, usd, inr, seg, mix, col) in enumerate(pricing_rows):
    y = 4.72 + ri * 0.58
    row_fill = WHITE if ri % 2 == 0 else LGRAY
    row_data = [plan, usd, inr, seg, mix]
    for ci, (val, cx, cw) in enumerate(zip(row_data, col_xs, col_ws)):
        rect(s, cx, y, cw - 0.06, 0.52, row_fill, BORDER, 0.5)
        tc = col if ci == 0 else (SUBTEXT if ci > 2 else TEXT)
        txt(s, val, cx + 0.08, y + 0.1, cw - 0.18, 0.35, size=10, bold=(ci == 0), color=tc)

    # Colored left strip
    rect(s, 0.45, y, 0.06, 0.52, col)

slide_num(s, 3)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 4 — REVENUE PROJECTIONS (Year 1)
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
hdr(s, "Revenue Projections — Year 1", "Month-by-month MRR forecast: conservative growth model")

monthly = [
    ("M1",  0,    0,    "Pre-launch"),
    ("M2",  3,    450,  "Early access"),
    ("M3",  8,    1200, "Public launch"),
    ("M4",  15,   2250, "Grow"),
    ("M5",  25,   3750, "Grow"),
    ("M6",  38,   5700, "Scale"),
    ("M7",  55,   8250, "Scale"),
    ("M8",  72,   10800,"Scale"),
    ("M9",  90,   13500,"Optimize"),
    ("M10", 108,  16200,"Optimize"),
    ("M11", 125,  18750,"Steady"),
    ("M12", 145,  21750,"End Y1"),
]

max_mrr = 21750
chart_h = 3.8
chart_y = 1.55
chart_x = 0.45
bar_area_w = 9.8
bar_w = bar_area_w / len(monthly) - 0.06
bg_chart = card(s, chart_x, chart_y, bar_area_w, chart_h, fill=WHITE, border=BORDER)

for i, (month, customers, mrr, note) in enumerate(monthly):
    bx = chart_x + i * (bar_w + 0.06) + 0.08
    if mrr > 0:
        bh = (mrr / max_mrr) * (chart_h - 0.55)
        by = chart_y + chart_h - bh - 0.3
        col = TEAL if mrr > 10000 else (LTEAL if mrr > 3000 else RGBColor(0xA7, 0xF3, 0xD0))
        bar = rect(s, bx, by, bar_w, bh, col)
        bar.line.fill.background()
        if mrr >= 5700:
            label = f"${mrr//1000}K"
            txt(s, label, bx, by - 0.28, bar_w, 0.25, size=7.5, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    txt(s, month, bx, chart_y + chart_h - 0.28, bar_w, 0.25, size=8, color=SUBTEXT, align=PP_ALIGN.CENTER)

# Y-axis labels
for v in [0, 5000, 10000, 15000, 20000]:
    vy = chart_y + chart_h - 0.3 - (v / max_mrr) * (chart_h - 0.55)
    txt(s, f"${v//1000}K", chart_x + bar_area_w + 0.08, vy - 0.14, 0.7, 0.3, size=8, color=MUTED)
    rect(s, chart_x + 0.05, vy, bar_area_w - 0.1, 0.01, BORDER)

txt(s, "Monthly Recurring Revenue (MRR) — USD", chart_x, 5.4, bar_area_w, 0.3, size=9, color=SUBTEXT, align=PP_ALIGN.CENTER)

# Right panel — summary stats
card(s, 10.45, 1.55, 2.63, 4.5, fill=NAVY_BG, border=NAVYMD, bw=1.0)
txt(s, "Year 1 Summary", 10.6, 1.68, 2.33, 0.35, size=11, bold=True, color=NAVY)
y1_stats = [
    ("End-MRR",   "$21.7K"),
    ("End-ARR",   "$261K"),
    ("Customers", "145"),
    ("Avg ARPU",  "$150/mo"),
    ("Churn",     "<5%"),
    ("Trial→Paid","35%"),
]
for i, (lbl, val) in enumerate(y1_stats):
    ys = 2.1 + i * 0.62
    txt(s, lbl, 10.6, ys, 1.1, 0.35, size=9.5, color=SUBTEXT)
    txt(s, val, 11.72, ys, 1.2, 0.35, size=10.5, bold=True, color=NAVY, align=PP_ALIGN.RIGHT)
    rect(s, 10.6, ys + 0.42, 2.33, 0.02, BORDER)

# 3-year outlook
txt(s, "3-Year Growth Outlook", 0.45, 5.78, 12.43, 0.35, size=12, bold=True, color=NAVY)
years = [
    ("Year 1", "145 customers", "$261K ARR", MUTED),
    ("Year 2", "500 customers", "$1.2M ARR", NAVYMD),
    ("Year 3", "1,500 customers","$5.4M ARR", TEAL),
]
for i, (yr, cust, arr, col) in enumerate(years):
    x3 = 0.45 + i * 4.3
    card(s, x3, 6.2, 4.05, 0.98, fill=WHITE, border=col, bw=1.5)
    rect(s, x3, 6.2, 4.05, 0.06, col)
    txt(s, yr, x3 + 0.15, 6.3, 1.3, 0.32, size=11, bold=True, color=col)
    txt(s, cust, x3 + 1.55, 6.3, 1.4, 0.32, size=10.5, color=SUBTEXT)
    txt(s, arr, x3 + 3.0, 6.3, 0.9, 0.32, size=11, bold=True, color=col, align=PP_ALIGN.RIGHT)

slide_num(s, 4)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 5 — GO-TO-MARKET STRATEGY
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
hdr(s, "Go-to-Market Strategy", "Multi-channel approach focused on high-intent AWS professionals")

channels = [
    ("LinkedIn Outreach",   "Direct outreach to AWS consultants, CTOs, and DevOps leads.\nPersonalized connection messages referencing their role.",
     "Primary", NAVY, "High intent, direct decision-makers"),
    ("Content Marketing",   "LinkedIn posts, YouTube tutorials about AWS security.\n'Top 5 AWS misconfigs' type content drives organic traffic.",
     "Primary", TEAL, "Builds authority, drives inbound"),
    ("ProductHunt / HN",   "Coordinated launches on ProductHunt and Hacker News.\nCommunity validation and early adopter acquisition.",
     "Launch", NAVYMD, "1,000+ signups in launch week"),
    ("SEO / Blog",          "Rank for 'AWS security scanner', 'CIS benchmark check'.\nTechnical blog posts attract high-intent searchers.",
     "Growth", AMBER, "Compounding, low CAC"),
    ("Consultant Referrals","AWS consultants refer client accounts, expand naturally.\nReferral incentive: 1 free month per paid referral.",
     "Growth", ORANGE, "Built-in network effect"),
    ("Cold Email Campaigns", "Targeted outreach to CTOs of funded startups needing SOC2.\nPersonalized, data-driven subject lines.",
     "Scale", PURPLE, "High conversion on SOC2 angle"),
]
for i, (chan, desc, phase, col, result) in enumerate(channels):
    ci = i % 3; row = i // 3
    x = 0.45 + ci * 4.3
    y = 1.55 + row * 2.7
    card(s, x, y, 4.05, 2.5, fill=WHITE, border=col, bw=1.0)
    rect(s, x, y, 4.05, 0.3, col)
    txt(s, f"{phase.upper()} CHANNEL", x + 0.15, y + 0.05, 3.75, 0.22, size=7.5, bold=True, color=WHITE)
    txt(s, chan, x + 0.18, y + 0.4, 3.7, 0.42, size=12.5, bold=True, color=TEXT)
    txt(s, desc, x + 0.18, y + 0.88, 3.7, 0.98, size=10, color=SUBTEXT)
    rect(s, x + 0.18, y + 2.1, 3.7, 0.28, col)
    txt(s, result, x + 0.25, y + 2.13, 3.55, 0.22, size=8.5, bold=True, color=WHITE)

slide_num(s, 5)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 6 — TARGET CUSTOMER SEGMENTS
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
hdr(s, "Target Customer Segments", "4 distinct, high-value customer profiles with clear pain points")

segments = [
    ("AWS Consultants",    "5–20 client accounts to manage.\nManual security reviews take 2 days each.",
     "VigiliCloud cuts review to 2 minutes.\nDeliver polished reports to every client.",
     "Pro / MSP", "$299–$999/mo", TEAL, TEAL_BG,
     "Scalability — more clients, same hours"),
    ("Startup CTOs",       "Accumulated 2–3 years of AWS config.\nNeed SOC2 for their next enterprise deal.",
     "Find and fix issues before the auditor.\nOne-click compliance evidence export.",
     "Starter / Pro", "$99–$299/mo", NAVYMD, NAVY_BG,
     "Compliance — close enterprise deals"),
    ("DevOps Teams",       "Internal AWS, no dedicated security team.\nSecurity Hub is complex to configure.",
     "Scheduled daily scans, AI summaries.\nFixes in minutes, not weeks.",
     "Pro", "$299/mo", ORANGE, ORG_BG,
     "Efficiency — security without overhead"),
    ("MSPs & Agencies",    "Manage cloud infra for 10+ clients.\nManual security doesn't scale.",
     "One platform, all clients.\nWhite-label ready reports for resale.",
     "MSP", "$999/mo", PURPLE, PURP_BG,
     "Scale — offer security as a service"),
]
for i, (title, pain, value, plan, price, col, bg_c, hook) in enumerate(segments):
    x = 0.45 + (i % 2) * 6.44
    y = 1.55 + (i // 2) * 2.8
    card(s, x, y, 6.14, 2.6, fill=bg_c, border=col, bw=1.5)

    rect(s, x, y, 6.14, 0.07, col)
    txt(s, title, x + 0.18, y + 0.18, 4.0, 0.38, size=14, bold=True, color=TEXT)

    price_box = rect(s, x + 4.7, y + 0.15, 1.25, 0.45, col)
    price_box.line.fill.background()
    txt(s, price, x + 4.7, y + 0.18, 1.25, 0.38, size=8.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    txt(s, "Pain:", x + 0.18, y + 0.65, 0.55, 0.28, size=9, bold=True, color=RED)
    txt(s, pain, x + 0.75, y + 0.65, 5.18, 0.5, size=9.5, color=SUBTEXT)
    txt(s, "Value:", x + 0.18, y + 1.2, 0.65, 0.28, size=9, bold=True, color=TEAL)
    txt(s, value, x + 0.88, y + 1.2, 5.06, 0.5, size=9.5, color=SUBTEXT)

    rect(s, x + 0.15, y + 1.85, 5.84, 0.6, col)
    txt(s, f"Key hook: {hook}", x + 0.28, y + 1.96, 5.6, 0.36, size=9.5, bold=True, color=WHITE)

slide_num(s, 6)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 7 — COMPETITOR ANALYSIS
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
hdr(s, "Competitor Analysis", "VigiliCloud occupies a unique, underserved position in the market")

# Positioning matrix label
txt(s, "Market Positioning Map (Price vs. Ease of Use)", 0.45, 1.52, 8.0, 0.35, size=12, bold=True, color=NAVY)

# Draw quadrant
quad_x = 0.45; quad_y = 1.92; quad_w = 7.8; quad_h = 5.2
card(s, quad_x, quad_y, quad_w, quad_h, fill=WHITE, border=BORDER)
# Quadrant lines
rect(s, quad_x + quad_w / 2 - 0.01, quad_y, 0.02, quad_h, BORDER)
rect(s, quad_x, quad_y + quad_h / 2 - 0.01, quad_w, 0.02, BORDER)
# Axis labels
txt(s, "Low Price", quad_x + 0.1, quad_y + quad_h / 2 - 0.15, 1.5, 0.3, size=9, color=MUTED, italic=True)
txt(s, "High Price", quad_x + quad_w - 1.7, quad_y + quad_h / 2 - 0.15, 1.5, 0.3, size=9, color=MUTED, italic=True)
txt(s, "Easy to Use", quad_x + quad_w / 2 - 1.1, quad_y + 0.1, 2.2, 0.28, size=9, color=MUTED, italic=True, align=PP_ALIGN.CENTER)
txt(s, "Complex / Expert", quad_x + quad_w / 2 - 1.3, quad_y + quad_h - 0.32, 2.6, 0.28, size=9, color=MUTED, italic=True, align=PP_ALIGN.CENTER)

# Competitors plotted as dots
competitors = [
    ("Lacework",     6.5, 1.5, ORANGE),   # high price, moderate complexity
    ("Wiz",          6.0, 2.0, RED),       # high price, somewhat easier
    ("Orca",         6.2, 1.8, PURPLE),
    ("AWS Sec Hub",  4.5, 3.2, AMBER),     # medium price, complex
    ("Prowler",      0.5, 3.8, SLATE),     # free, complex
    ("Manual Audit", 1.2, 3.5, MUTED),     # free (time), complex
]
for comp, px, py, col in competitors:
    dot_x = quad_x + (px / 8) * quad_w
    dot_y = quad_y + (py / 4.5) * quad_h
    circle = rect(s, dot_x - 0.12, dot_y - 0.12, 0.24, 0.24, col)
    circle.line.fill.background()
    txt(s, comp, dot_x - 0.8, dot_y - 0.42, 1.6, 0.28, size=8.5, color=col, align=PP_ALIGN.CENTER)

# VigiliCloud — highlight
vc_x = quad_x + (1.2 / 8) * quad_w
vc_y = quad_y + (0.6 / 4.5) * quad_h
vc_dot = rect(s, vc_x - 0.2, vc_y - 0.2, 0.4, 0.4, TEAL)
vc_dot.line.color.rgb = NAVY; vc_dot.line.width = Pt(1.5)
txt(s, "VigiliCloud", vc_x - 0.85, vc_y - 0.5, 1.7, 0.28, size=9.5, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

# Right: comparison table
txt(s, "Head-to-Head Comparison", 8.45, 1.52, 4.73, 0.35, size=12, bold=True, color=NAVY)
comp_rows = [
    ("Company",      "Price",     "Setup", "AI",  "SMB?"),
    ("Wiz",          "$50K+/yr",  "Weeks", "Yes", "No"),
    ("Lacework",     "$30K+/yr",  "Days",  "Some","No"),
    ("Orca",         "$20K+/yr",  "Days",  "Some","No"),
    ("AWS Sec Hub",  "Variable",  "Hours", "No",  "Partial"),
    ("Prowler",      "Free",      "Dev",   "No",  "Dev-only"),
    ("VigiliCloud",  "$99/mo",    "5 min", "Yes", "Yes"),
]
col_xs2 = [8.45, 9.7, 10.75, 11.55, 12.2]
col_ws2 = [1.2,  1.0,  0.75,  0.6,   0.9]
for ri, row in enumerate(comp_rows):
    ry = 1.95 + ri * 0.58
    is_header = ri == 0
    is_vc = ri == len(comp_rows) - 1
    fill = NAVY if is_header else (TEAL_BG if is_vc else (WHITE if ri % 2 == 1 else LGRAY))
    for ci, (cell, cx, cw) in enumerate(zip(row, col_xs2, col_ws2)):
        rect(s, cx, ry, cw - 0.04, 0.52, fill, BORDER if not is_vc else TEAL, 0.5)
        tc = WHITE if is_header else (TEAL if is_vc else (NAVY if ci == 0 else SUBTEXT))
        txt(s, cell, cx + 0.06, ry + 0.1, cw - 0.14, 0.34, size=9,
            bold=(is_header or is_vc or ci == 0), color=tc,
            align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER)

slide_num(s, 7)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 8 — PRODUCT ROADMAP
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
hdr(s, "Product Roadmap", "Clear path from SMB AWS scanner to multi-cloud security platform")

# Timeline bar
rect(s, 0.45, 2.22, 12.43, 0.06, NAVY)

quarters = [
    ("Q2 2026\n(Now)",         [
        "Multi-account dashboard",
        "Scheduled daily scans",
        "Razorpay billing live",
        "Email alerts system",
    ], TEAL, 0.45),
    ("Q3 2026",                [
        "Slack / Teams alerts",
        "API access for devs",
        "Custom scan policies",
        "Bulk account import",
    ], NAVYMD, 3.6),
    ("Q4 2026",                [
        "Azure & GCP support",
        "GitHub Actions CI/CD",
        "HIPAA / PCI checks",
        "White-label MSP mode",
    ], ORANGE, 6.75),
    ("2027+",                  [
        "AI auto-remediation",
        "IaC scanning (Terraform)",
        "SOC2 Type II evidence",
        "Enterprise SIEM export",
    ], PURPLE, 9.9),
]
for q_title, items, col, x in quarters:
    # Dot on timeline
    rect(s, x + 1.4, 2.1, 0.24, 0.24, col)

    # Quarter label
    txt(s, q_title, x, 1.55, 3.0, 0.52, size=11, bold=True, color=col, align=PP_ALIGN.CENTER)

    # Card
    card(s, x, 2.45, 3.05, 4.55, fill=WHITE, border=col, bw=1.5)
    rect(s, x, 2.45, 3.05, 0.08, col)

    for j, item in enumerate(items):
        iy = 2.62 + j * 0.95
        card(s, x + 0.15, iy, 2.75, 0.8, fill=LGRAY, border=BORDER, bw=0.5)
        rect(s, x + 0.15, iy, 0.05, 0.8, col)
        txt(s, item, x + 0.3, iy + 0.2, 2.45, 0.42, size=10, color=TEXT)

txt(s, "Roadmap is directional — priorities adjusted based on customer feedback and market response",
    0.45, 7.12, 12.43, 0.28, size=9.5, color=MUTED, italic=True, align=PP_ALIGN.CENTER)

slide_num(s, 8)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 9 — CURRENT TRACTION & MILESTONES
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
hdr(s, "Traction & Milestones", "What's been built and what's already working")

# What's live
txt(s, "Product — Live & Operational", 0.45, 1.52, 6.2, 0.38, size=13, bold=True, color=NAVY)

live_items = [
    ("Full-stack SaaS deployed",      "Next.js UI + FastAPI backend on Render, live at app.vigilicloud.com", TEAL),
    ("10 security checks running",    "S3, IAM, EC2, RDS, VPC, KMS, EBS, MFA, CloudTrail, VPC Flow Logs", TEAL),
    ("AI analysis integrated",        "Claude AI generates plain-English executive summaries per scan", TEAL),
    ("IAM role-based auth",           "Read-only STS-based AWS access, zero credential storage", TEAL),
    ("Scan results & exports",        "CSV/JSON finding exports, severity filtering, per-resource findings", TEAL),
    ("Razorpay billing ready",        "Payment system integrated, INR + USD pricing, subscription management", AMBER),
    ("Email alert system",            "CRITICAL finding notifications via Resend email API", TEAL),
]
for i, (title, desc, col) in enumerate(live_items):
    yi = 1.98 + i * 0.6
    status_box = rect(s, 0.45, yi + 0.08, 0.55, 0.38, TEAL_BG if col == TEAL else AMB_BG, col, 0.75)
    txt(s, "LIVE" if col == TEAL else "SOON", 0.45, yi + 0.1, 0.55, 0.34, size=7.5, bold=True,
        color=col, align=PP_ALIGN.CENTER)
    txt(s, title, 1.12, yi + 0.04, 5.05, 0.3, size=11, bold=True, color=TEXT)
    txt(s, desc, 1.12, yi + 0.34, 5.05, 0.25, size=9, color=SUBTEXT)

# Right: metrics / milestones
txt(s, "Business Milestones", 7.0, 1.52, 6.18, 0.38, size=13, bold=True, color=NAVY)

milestones = [
    ("Product launched",        "app.vigilicloud.com is live", TEAL),
    ("Training deck created",   "Client pitch + product training complete", TEAL),
    ("Razorpay verified",       "Payment processing ready to activate", AMBER),
    ("Domain email setup",      "vigilicloud.com transactional email pending", AMBER),
    ("LinkedIn campaign start", "Outreach to AWS consultants planned", NAVY),
    ("ProductHunt launch",      "Coordinated launch scheduled for Q2 2026", NAVY),
    ("First 10 paying clients", "Target: end of Q2 2026", NAVY),
    ("$1K MRR milestone",       "Target: Month 5 of active GTM", NAVY),
]
for i, (title, desc, col) in enumerate(milestones):
    yi = 1.98 + i * 0.66
    if col == TEAL:
        status = "DONE"
        bg_c = TEAL_BG
    elif col == AMBER:
        status = "IN PROG"
        bg_c = AMB_BG
    else:
        status = "PLANNED"
        bg_c = NAVY_BG

    card(s, 7.0, yi, 6.18, 0.55, fill=bg_c, border=col, bw=0.5)
    rect(s, 7.0, yi, 0.06, 0.55, col)
    status_b = rect(s, 7.08, yi + 0.12, 0.88, 0.3, col)
    status_b.line.fill.background()
    txt(s, status, 7.08, yi + 0.13, 0.88, 0.26, size=7, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, title, 8.04, yi + 0.04, 3.8, 0.26, size=10, bold=True, color=TEXT)
    txt(s, desc, 8.04, yi + 0.3, 3.8, 0.22, size=8.5, color=SUBTEXT)

slide_num(s, 9)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 10 — MARKETING TARGETS & KPIs
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
hdr(s, "Marketing Targets & KPIs", "Clear, measurable goals for the next 12 months")

# Funnel metrics
txt(s, "Growth Funnel — 12-Month Targets", 0.45, 1.52, 8.0, 0.35, size=13, bold=True, color=NAVY)

funnel = [
    ("Website Visitors",   "5,000/mo",   "via SEO + content", MUTED),
    ("Trial Signups",      "500/mo",     "10% of visitors",   SLATE),
    ("Activated Accounts", "175/mo",     "35% of signups",    NAVYMD),
    ("Paying Customers",   "52/mo",      "30% of activated",  TEAL),
    ("MRR (Month 12)",     "$21,750",    "145 customers avg", DKGREEN),
]
for i, (stage, target, note, col) in enumerate(funnel):
    bar_w = 8.0 - i * 1.4
    bx = 0.45 + (8.0 - bar_w) / 2
    fy = 1.98 + i * 0.88
    bar = rect(s, bx, fy, bar_w, 0.62, col)
    bar.line.fill.background()
    txt(s, stage, bx + 0.2, fy + 0.1, 2.5, 0.42, size=10.5, bold=True, color=WHITE)
    txt(s, target, bx + bar_w - 1.5, fy + 0.1, 1.3, 0.42, size=11, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
    txt(s, note, bx + bar_w + 0.15, fy + 0.18, 2.5, 0.3, size=9.5, color=SUBTEXT)

# Right: quarterly OKRs
txt(s, "Quarterly OKRs", 8.6, 1.52, 4.58, 0.35, size=13, bold=True, color=NAVY)

okrs = [
    ("Q2 2026", [
        "Launch publicly on ProductHunt",
        "First 10 paying customers",
        "Reach $1K MRR",
    ], TEAL),
    ("Q3 2026", [
        "50 paying customers",
        "Reach $7K MRR",
        "Launch Slack integration",
    ], NAVYMD),
    ("Q4 2026", [
        "100 customers, $15K MRR",
        "Launch Azure support",
        "First MSP partner signed",
    ], ORANGE),
]
for i, (quarter, goals, col) in enumerate(okrs):
    y_okr = 1.98 + i * 1.68
    card(s, 8.6, y_okr, 4.58, 1.55, fill=WHITE, border=col, bw=1.5)
    rect(s, 8.6, y_okr, 4.58, 0.08, col)
    txt(s, quarter, 8.75, y_okr + 0.14, 4.25, 0.3, size=11, bold=True, color=col)
    for j, goal in enumerate(goals):
        rect(s, 8.75, y_okr + 0.5 + j * 0.33, 0.14, 0.22, col)
        txt(s, goal, 8.97, y_okr + 0.47 + j * 0.33, 3.98, 0.28, size=9.5, color=SUBTEXT)

slide_num(s, 10)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 11 — COMPETITOR REVENUE CONTEXT
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
hdr(s, "Industry Revenue Context", "Where VigiliCloud fits in the cloud security revenue landscape")

comp_rev = [
    ("Wiz",           "2020", "$500M ARR", "Cloud-native security platform", ORANGE),
    ("Lacework",      "2015", "$150M ARR", "Behavioral security analytics",  RED),
    ("Orca Security", "2019", "$100M ARR", "Cloud risk management",          PURPLE),
    ("Snyk",          "2015", "$300M ARR", "Developer security platform",    NAVYMD),
    ("Prowler",       "2016", "Open Source","CLI-based AWS audit tool",       SLATE),
    ("VigiliCloud",   "2026", "Growing",   "SMB AWS security scanning",      TEAL),
]
for i, (company, founded, revenue, desc, col) in enumerate(comp_rev):
    x = 0.45 + (i % 3) * 4.3
    y = 1.55 + (i // 3) * 2.65
    is_vc = company == "VigiliCloud"
    fill = TEAL_BG if is_vc else WHITE
    bw = 2.0 if is_vc else 0.75
    card(s, x, y, 4.05, 2.4, fill=fill, border=col, bw=bw)
    if is_vc:
        rect(s, x, y, 4.05, 0.08, TEAL)

    txt(s, company, x + 0.18, y + 0.18, 2.8, 0.42, size=14, bold=True, color=col)
    txt(s, f"Founded {founded}", x + 0.18, y + 0.62, 1.8, 0.28, size=9.5, color=MUTED)
    rev_box = rect(s, x + 2.45, y + 0.55, 1.42, 0.38, col)
    rev_box.line.fill.background()
    txt(s, revenue, x + 2.45, y + 0.58, 1.42, 0.32, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, desc, x + 0.18, y + 1.0, 3.7, 0.4, size=10, color=SUBTEXT)

    if is_vc:
        txt(s, "YOUR OPPORTUNITY: The SMB gap that none of these players fill",
            x + 0.18, y + 1.5, 3.7, 0.62, size=9.5, bold=True, color=TEAL)

slide_num(s, 11)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 12 — FUTURE UPDATES VISION
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
hdr(s, "Product Vision — Where We're Going", "From AWS scanner to the #1 cloud security platform for SMBs")

vision_phases = [
    ("Phase 1\nToday",        "AWS Security Scanner",
     ["10 CIS checks", "AI analysis", "Compliance exports", "Fix guidance"],
     "Complete",    TEAL,   TEAL_BG),
    ("Phase 2\nQ3–Q4 2026",   "Multi-Cloud Platform",
     ["Azure + GCP support", "Slack/Teams alerts", "CI/CD integration", "API access"],
     "In Progress", NAVYMD, NAVY_BG),
    ("Phase 3\n2027",         "AI Remediation Engine",
     ["Auto-fix suggestions", "IaC scanning", "HIPAA/PCI modules", "SOC2 evidence"],
     "Planned",     ORANGE, ORG_BG),
    ("Phase 4\n2027+",        "Enterprise Platform",
     ["White-label MSP", "Custom frameworks", "SIEM integration", "24/7 monitoring"],
     "Vision",      PURPLE, PURP_BG),
]
for i, (phase, title, items, status, col, bg_c) in enumerate(vision_phases):
    x = 0.45 + i * 3.22
    card(s, x, 1.55, 3.0, 5.2, fill=bg_c, border=col, bw=1.5)

    status_b = rect(s, x + 0.18, 1.72, 2.64, 0.35, col)
    status_b.line.fill.background()
    txt(s, status.upper(), x + 0.18, 1.74, 2.64, 0.3, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    txt(s, phase, x + 0.18, 2.14, 2.64, 0.55, size=11, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(s, title, x + 0.18, 2.74, 2.64, 0.55, size=12.5, bold=True, color=TEXT, align=PP_ALIGN.CENTER)

    rect(s, x + 0.35, 3.35, 2.3, 0.03, col)

    for j, item in enumerate(items):
        iy = 3.48 + j * 0.68
        rect(s, x + 0.28, iy + 0.1, 0.22, 0.22, col)
        txt(s, item, x + 0.62, iy, 2.22, 0.58, size=10.5, color=SUBTEXT)

# Arrow connectors between phases
for i in range(3):
    ax = 0.45 + (i + 1) * 3.22 - 0.18
    txt(s, "→", ax, 3.5, 0.25, 0.4, size=18, bold=True, color=MUTED, align=PP_ALIGN.CENTER)

txt(s, "Each phase expands the moat — switching costs grow as customers integrate deeper compliance workflows",
    0.45, 6.93, 12.43, 0.35, size=10.5, color=SUBTEXT, align=PP_ALIGN.CENTER, italic=True)

slide_num(s, 12)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 13 — CALL TO ACTION / NEXT STEPS
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, NAVY)
rect(s, 0, 0, 13.33, 0.08, GOLD)
rect(s, 0, 7.42, 13.33, 0.08, TEAL)

txt(s, "VigiliCloud is solving a real problem", 0.8, 0.4, 11.73, 0.7,
    size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "in a massive market — with a product that's live and working today.",
    0.8, 1.12, 11.73, 0.5, size=18, color=LTEAL, align=PP_ALIGN.CENTER)

rect(s, 4.0, 1.78, 5.33, 0.05, GOLD)

# 3 action cards
actions = [
    ("Try It Free",      "app.vigilicloud.com\n2-week trial, no credit card",    TEAL,   "For potential customers"),
    ("Partner With Us",  "MSP & referral partnerships\nleelakrishnakoppolu@gmail.com", NAVYMD, "For agencies & consultants"),
    ("Give Feedback",    "Book a 20-min call\ncalendly.com/leelakrishnakoppolu", GOLD,   "For advisors & investors"),
]
for i, (action, detail, col, desc) in enumerate(actions):
    x = 0.65 + i * 4.05
    act_card = rect(s, x, 2.05, 3.75, 2.5, RGBColor(0x1a, 0x3a, 0x7a), col, 2.0)
    rect(s, x, 2.05, 3.75, 0.07, col)
    txt(s, action, x + 0.2, 2.2, 3.35, 0.45, size=15, bold=True, color=col)
    txt(s, detail, x + 0.2, 2.72, 3.35, 0.72, size=11, color=RGBColor(0xBF, 0xDB, 0xFF))
    desc_b = rect(s, x + 0.2, 4.1, 3.35, 0.32, col)
    desc_b.line.fill.background()
    txt(s, desc, x + 0.2, 4.12, 3.35, 0.28, size=8.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Stats reminder
stats_row = [
    ("$50B", "Market Size"),
    ("78%", "Accounts with CRITICAL issues"),
    ("$99/mo", "Starting Price"),
    ("2 min", "Full AWS Scan"),
    ("0", "Stored Credentials"),
]
for i, (val, lbl) in enumerate(stats_row):
    sx = 0.5 + i * 2.48
    txt(s, val, sx, 5.0, 2.28, 0.55, size=22, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    txt(s, lbl, sx, 5.55, 2.28, 0.38, size=9, color=MUTED, align=PP_ALIGN.CENTER)

rect(s, 0.45, 5.98, 12.43, 0.04, RGBColor(0x2a, 0x4a, 0x9a))

txt(s, "Leela Krishna Koppolu  ·  leelakrishnakoppolu@gmail.com  ·  app.vigilicloud.com  ·  May 2026",
    0.8, 6.12, 11.73, 0.38, size=11, color=MUTED, align=PP_ALIGN.CENTER)
txt(s, "VigiliCloud — AWS Security for Every Business. Not Just Enterprise.",
    0.8, 6.58, 11.73, 0.42, size=13, bold=True, color=LTEAL, align=PP_ALIGN.CENTER)

slide_num(s, 13)

# ── Save ──────────────────────────────────────────────────────────
out = r"c:\Users\leela\compliance-ai-saas\VigiliCloud_Marketing_Deck.pptx"
prs.save(out)
print(f"Saved: {out}")
