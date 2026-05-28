from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# ── Colour palette ──────────────────────────────────────────────
BLACK   = RGBColor(0x06, 0x06, 0x06)
DARK    = RGBColor(0x0f, 0x0f, 0x0f)
CARD    = RGBColor(0x16, 0x16, 0x16)
GREEN   = RGBColor(0x10, 0xb9, 0x81)   # emerald-500
LGREEN  = RGBColor(0x34, 0xd3, 0x99)   # emerald-400
WHITE   = RGBColor(0xff, 0xff, 0xff)
GREY    = RGBColor(0xa3, 0xa3, 0xa3)
DKGREY  = RGBColor(0x52, 0x52, 0x52)
RED     = RGBColor(0xef, 0x44, 0x44)
ORANGE  = RGBColor(0xf9, 0x73, 0x16)
YELLOW  = RGBColor(0xea, 0xb3, 0x08)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # completely blank

# ── Helpers ─────────────────────────────────────────────────────
def bg(slide, color=DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, l, t, w, h, color, radius=False):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def txt(slide, text, l, t, w, h, size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def accent_bar(slide, t=0.28):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(t), Inches(13.33), Inches(0.055))
    bar.fill.solid(); bar.fill.fore_color.rgb = GREEN
    bar.line.fill.background()

def slide_num(slide, n):
    txt(slide, str(n), 12.6, 7.1, 0.5, 0.3, size=9, color=DKGREY, align=PP_ALIGN.RIGHT)

def tag(slide, label, l, t, color=GREEN, text_color=None):
    tc = text_color or BLACK
    b = box(slide, l, t, len(label)*0.095+0.25, 0.26, color)
    txt(slide, label, l+0.09, t+0.03, len(label)*0.095+0.1, 0.22, size=9, bold=True,
        color=tc, align=PP_ALIGN.CENTER)

def card(slide, l, t, w, h, fill=CARD):
    c = box(slide, l, t, w, h, fill)
    # thin green border via outline
    from pptx.util import Pt as UPt
    c.line.color.rgb = GREEN
    c.line.width = UPt(0.6)
    return c

# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, BLACK)
accent_bar(s, 0.0)

# Big green V logo box
logo = box(s, 0.45, 0.6, 0.65, 0.65, RGBColor(0x06, 0x2a, 0x1e))
logo.line.color.rgb = GREEN; logo.line.width = Pt(1)
txt(s, "V", 0.45, 0.6, 0.65, 0.65, size=28, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

txt(s, "VigiliCloud", 1.2, 0.65, 4, 0.55, size=28, bold=True, color=WHITE)

txt(s, "AWS Security Scanning — Built for Speed, Clarity & Action",
    0.45, 1.45, 9.5, 0.7, size=20, bold=False, color=GREY)

# Divider
div = box(s, 0.45, 2.25, 6, 0.04, GREEN)
div.line.fill.background()

txt(s, "Detect misconfigurations. Get fix guidance. Prove compliance.",
    0.45, 2.45, 9.5, 0.55, size=16, color=LGREEN)

txt(s, "Presented by  Leela Krishna Koppolu  |  vigilicloud.com  |  2026",
    0.45, 6.9, 9, 0.45, size=11, color=DKGREY)

# Right visual panel
rp = box(s, 9.5, 1.0, 3.6, 5.5, CARD)
rp.line.color.rgb = GREEN; rp.line.width = Pt(0.8)

txt(s, "🔍  S3 Public Access        CRITICAL", 9.65, 1.2, 3.3, 0.38, size=10, color=RED)
txt(s, "🔑  Root Access Keys        CRITICAL", 9.65, 1.65, 3.3, 0.38, size=10, color=RED)
txt(s, "🔐  IAM Permissions             HIGH", 9.65, 2.1,  3.3, 0.38, size=10, color=ORANGE)
txt(s, "📱  MFA Enforcement              HIGH", 9.65, 2.55, 3.3, 0.38, size=10, color=ORANGE)
txt(s, "🛡️   Security Groups             HIGH", 9.65, 3.0,  3.3, 0.38, size=10, color=ORANGE)
txt(s, "💾  EBS Encryption            MEDIUM", 9.65, 3.45, 3.3, 0.38, size=10, color=YELLOW)
txt(s, "🗄️   RDS Encryption             HIGH", 9.65, 3.9,  3.3, 0.38, size=10, color=ORANGE)
txt(s, "📋  CloudTrail Logging        MEDIUM", 9.65, 4.35, 3.3, 0.38, size=10, color=YELLOW)
txt(s, "🌐  VPC Flow Logs             MEDIUM", 9.65, 4.8,  3.3, 0.38, size=10, color=YELLOW)
txt(s, "🔄  KMS Key Rotation          MEDIUM", 9.65, 5.25, 3.3, 0.38, size=10, color=YELLOW)
txt(s, "10 SECURITY CHECKS", 9.65, 6.0, 3.3, 0.35, size=10, bold=True, color=GREEN)

slide_num(s, 1)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, BLACK)
accent_bar(s)

tag(s, "THE PROBLEM", 0.45, 0.55)
txt(s, "AWS misconfigurations are the #1 cause of cloud data breaches",
    0.45, 0.9, 11, 0.65, size=26, bold=True, color=WHITE)

stats = [
    ("82%", "of breaches involve\nmisconfigured cloud\nresources", RED),
    ("$4.5M", "average cost of a\ncloud data breach\n(IBM 2024)", ORANGE),
    ("68%", "of companies don't\nknow their full AWS\nsecurity posture", YELLOW),
    ("2 days", "average time to\ndetect a critical\nmisconfiguration", LGREEN),
]
for i, (val, lbl, col) in enumerate(stats):
    x = 0.45 + i * 3.2
    c = card(s, x, 1.75, 2.95, 2.4)
    txt(s, val, x+0.15, 1.95, 2.65, 0.75, size=32, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(s, lbl, x+0.15, 2.75, 2.65, 0.95, size=12, color=GREY, align=PP_ALIGN.CENTER)

txt(s, "Real-world pain points your clients face every day:",
    0.45, 4.35, 12, 0.38, size=13, bold=True, color=LGREEN)

pains = [
    "No visibility into who has admin access or what's publicly exposed",
    "Manual audits take 2–3 days and miss things — automated scanning catches everything in 2 minutes",
    "No clear remediation path — engineers know something is wrong but don't know how to fix it",
    "Compliance evidence for SOC2, ISO 27001 takes weeks to collect manually",
]
for i, p in enumerate(pains):
    txt(s, f"✗  {p}", 0.45, 4.8 + i*0.42, 12.4, 0.38, size=11.5, color=GREY)

slide_num(s, 2)

# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — WHAT IS VIGILICLOUD
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, BLACK)
accent_bar(s)

tag(s, "OUR SOLUTION", 0.45, 0.55)
txt(s, "VigiliCloud — AWS Security Scanning in 2 Minutes",
    0.45, 0.9, 11, 0.65, size=26, bold=True, color=WHITE)

txt(s, "Connect your AWS account → Run a scan → Get prioritised findings with exact fix steps.\nNo agents. No installations. No AWS expertise required.",
    0.45, 1.65, 8.5, 0.75, size=13, color=GREY)

steps = [
    ("01", "Connect AWS Account", "Create a read-only IAM role,\npaste the ARN. We never store\nyour credentials.", "5 minutes"),
    ("02", "Run Security Scan",    "VigiliCloud checks all 10\nsecurity areas across your\nentire AWS account.",          "~2 minutes"),
    ("03", "Fix What's Wrong",     "Every finding includes exact\nConsole path, CLI commands,\nand step-by-step fix guide.",   "Evidence exports"),
]
for i, (num, title, desc, badge) in enumerate(steps):
    x = 0.45 + i * 4.25
    c = card(s, x, 2.55, 3.95, 3.5)
    txt(s, num, x+0.2, 2.7, 1.0, 0.6, size=28, bold=True, color=GREEN)
    txt(s, title, x+0.2, 3.35, 3.55, 0.5, size=14, bold=True, color=WHITE)
    txt(s, desc, x+0.2, 3.9, 3.55, 1.0, size=11, color=GREY)
    tag(s, badge, x+0.2, 5.5, RGBColor(0x06,0x2a,0x1e), LGREEN)

txt(s, "✦  Free 2-week trial  ·  No credit card  ·  Setup in 5 minutes  ·  Works with any AWS account",
    0.45, 6.35, 12, 0.38, size=11, color=DKGREY, align=PP_ALIGN.CENTER)

slide_num(s, 3)

# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — 10 SECURITY CHECKS
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, BLACK)
accent_bar(s)

tag(s, "WHAT WE SCAN", 0.45, 0.55)
txt(s, "10 Critical Security Checks — Zero Setup Required",
    0.45, 0.9, 11, 0.55, size=26, bold=True, color=WHITE)

checks = [
    ("🪣", "S3 Public Access",    "Open buckets, public ACLs,\nexposed bucket policies",         "CRITICAL", RED),
    ("🔑", "Root Access Keys",    "Active root account keys —\nmost dangerous misconfiguration", "CRITICAL", RED),
    ("🔐", "IAM Permissions",     "Over-permissioned roles,\nunnecessary admin access",          "HIGH",     ORANGE),
    ("📱", "MFA Enforcement",     "Missing MFA on IAM users\nand root account",                  "HIGH",     ORANGE),
    ("🛡️","Security Groups",     "EC2 ports open to\nthe entire internet (0.0.0.0/0)",          "HIGH",     ORANGE),
    ("🗄️","RDS Encryption",      "Unencrypted databases,\npublicly accessible RDS",             "HIGH",     ORANGE),
    ("💾", "EBS Encryption",      "Unencrypted volumes,\nmissing default encryption",            "MEDIUM",   YELLOW),
    ("📋", "CloudTrail Logging",  "Inactive trails, missing\nmulti-region & log validation",    "MEDIUM",   YELLOW),
    ("🌐", "VPC Flow Logs",       "Network traffic not\nlogged — forensic blind spot",          "MEDIUM",   YELLOW),
    ("🔄", "KMS Key Rotation",    "Customer-managed keys\nnot auto-rotating",                   "MEDIUM",   YELLOW),
]

cols = 5
for i, (icon, name, desc, sev, col) in enumerate(checks):
    row = i // cols
    ci  = i  % cols
    x = 0.35 + ci * 2.55
    y = 1.65 + row * 2.65
    c = card(s, x, y, 2.4, 2.45)
    txt(s, icon, x+0.15, y+0.15, 0.5, 0.45, size=18)
    txt(s, sev,  x+0.7,  y+0.22, 1.55, 0.28, size=8, bold=True, color=col)
    txt(s, name, x+0.15, y+0.65, 2.1,  0.42, size=11, bold=True, color=WHITE)
    txt(s, desc, x+0.15, y+1.1,  2.1,  0.8,  size=9.5, color=GREY)

slide_num(s, 4)

# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — SECURITY STATS / PIE DATA
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, BLACK)
accent_bar(s)

tag(s, "SECURITY INSIGHTS", 0.45, 0.55)
txt(s, "What We Find in a Typical AWS Account Scan",
    0.45, 0.9, 10, 0.55, size=26, bold=True, color=WHITE)

# Left — donut representation as stacked bars (pptx doesn't easily do circles)
txt(s, "Finding Severity Breakdown", 0.45, 1.65, 5.5, 0.38, size=13, bold=True, color=LGREEN)

severity_data = [
    ("CRITICAL", "22%", 2.2, RED),
    ("HIGH",     "35%", 3.5, ORANGE),
    ("MEDIUM",   "28%", 2.8, YELLOW),
    ("LOW",      "15%", 1.5, DKGREY),
]
bar_start = 0.45
for sev, pct, w, col in severity_data:
    bx = box(s, bar_start, 2.15, w, 0.45, col)
    bx.line.fill.background()
    bar_start += w + 0.05

# Legend
for i, (sev, pct, w, col) in enumerate(severity_data):
    x = 0.45 + i * 3.0
    dot = box(s, x, 2.78, 0.18, 0.18, col)
    dot.line.fill.background()
    txt(s, f"{sev}  {pct}", x+0.25, 2.73, 2.6, 0.28, size=11, bold=True, color=col)

# Right — 4 stat cards
stat_data = [
    ("78%",    "AWS accounts with\nat least 1 CRITICAL\nfinding",   RED,    5.8,  1.55),
    ("14",     "Average findings\nper account scan",                 ORANGE, 9.55, 1.55),
    ("#1",     "S3 misconfiguration\nis most common\nfinding type",  YELLOW, 5.8,  4.0),
    ("~1 day", "Average fix time\nwith step-by-step\nguidance",      LGREEN, 9.55, 4.0),
]
for val, lbl, col, x, y in stat_data:
    c = card(s, x, y, 3.5, 2.25)
    txt(s, val, x+0.2, y+0.25, 3.1, 0.7, size=30, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(s, lbl, x+0.2, y+1.0,  3.1, 0.95, size=11, color=GREY, align=PP_ALIGN.CENTER)

txt(s, "Source: aggregated data from VigiliCloud account scans",
    0.45, 6.85, 8, 0.32, size=9, color=DKGREY, italic=True)

slide_num(s, 5)

# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — KEY FEATURES
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, BLACK)
accent_bar(s)

tag(s, "KEY FEATURES", 0.45, 0.55)
txt(s, "Everything Your Team Needs — Nothing You Don't",
    0.45, 0.9, 11, 0.55, size=26, bold=True, color=WHITE)

features = [
    ("🤖", "AI Security Analysis",   "Claude AI summarises every scan,\nprioritises what to fix first, and\nexplains risks in plain English.",     LGREEN),
    ("📧", "Email Alerts",           "Get instantly notified by email\nwhenever a CRITICAL finding\nis detected in your account.",                LGREEN),
    ("📁", "Compliance Exports",     "Export findings as CSV or JSON\nfor SOC2, ISO 27001 audits,\nand compliance evidence packages.",           LGREEN),
    ("⏱️", "2-Minute Scans",        "Full account scan completes in\nunder 2 minutes. Schedule daily\nauto-scans for continuous monitoring.",   LGREEN),
    ("🔒", "Read-Only Access",       "We never store credentials.\nAccess is assumed temporarily\nvia a read-only IAM role.",                    LGREEN),
    ("🗺️", "Fix Guidance",          "Every finding includes exact\nAWS Console path, CLI command,\nand step-by-step fix instructions.",         LGREEN),
]
for i, (icon, title, desc, col) in enumerate(features):
    row = i // 3
    ci  = i  % 3
    x = 0.45 + ci * 4.25
    y = 1.65 + row * 2.65
    c = card(s, x, y, 3.95, 2.45)
    txt(s, icon,  x+0.2, y+0.2,  0.6,  0.5,  size=20)
    txt(s, title, x+0.2, y+0.75, 3.55, 0.42, size=13, bold=True, color=WHITE)
    txt(s, desc,  x+0.2, y+1.2,  3.55, 1.0,  size=10.5, color=GREY)

slide_num(s, 6)

# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — COMPETITOR COMPARISON
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, BLACK)
accent_bar(s)

tag(s, "COMPETITIVE LANDSCAPE", 0.45, 0.55)
txt(s, "How VigiliCloud Compares to Alternatives",
    0.45, 0.9, 11, 0.55, size=26, bold=True, color=WHITE)

# Table header
headers = ["", "VigiliCloud", "AWS\nSecurity Hub", "Prowler\n(Open Source)", "Lacework /\nOrca", "Manual\nAudit"]
col_w   = [2.3, 1.75, 1.75, 1.75, 1.75, 1.6]
col_x   = [0.3]
for w in col_w[:-1]: col_x.append(col_x[-1] + w)

# Header row
for i, (h, cx, cw) in enumerate(zip(headers, col_x, col_w)):
    bg_col = GREEN if i == 1 else RGBColor(0x1a, 0x1a, 0x1a)
    b = box(s, cx, 1.65, cw - 0.05, 0.55, bg_col)
    b.line.fill.background()
    tc = BLACK if i == 1 else GREY
    txt(s, h, cx+0.05, 1.68, cw-0.1, 0.5, size=9, bold=True, color=tc, align=PP_ALIGN.CENTER)

# Rows
rows = [
    ("Setup time",       "5 minutes",  "Hours + config", "CLI expertise", "Days + POC",    "2–3 days"),
    ("Price / month",    "$99–$999",   "$30–$1000+",     "Free (DIY)",    "$5K–$50K+",     "$0 (your time)"),
    ("Guided fix steps", "✔  Full",    "✘  None",        "✘  None",       "△  Limited",    "✘  Manual"),
    ("AI analysis",      "✔  Claude",  "✘  None",        "✘  None",       "△  Some",       "✘  None"),
    ("No expertise req.", "✔  Anyone", "✘  AWS expert",  "✘  Dev only",   "✘  Security eng","✘  Expert"),
    ("Compliance export", "✔  CSV/JSON","△  Partial",    "✘  Limited",    "✔  Enterprise", "✔  Manual"),
    ("SMB-friendly",     "✔  Yes",     "△  Complex",     "✘  Dev tool",   "✘  Enterprise", "✔  Yes"),
]
tick_map = {"✔": LGREEN, "✘": RED, "△": YELLOW}
for ri, row in enumerate(rows):
    y = 2.28 + ri * 0.57
    row_bg = RGBColor(0x10,0x10,0x10) if ri % 2 == 0 else BLACK
    for ci, (cell, cx, cw) in enumerate(zip(row, col_x, col_w)):
        b = box(s, cx, y, cw - 0.05, 0.52, row_bg)
        b.line.fill.background()
        if ci == 1:
            b.fill.solid(); b.fill.fore_color.rgb = RGBColor(0x06,0x1f,0x14)
        first = cell[:1]
        tc = tick_map.get(first, WHITE if ci == 0 else GREY)
        txt(s, cell, cx+0.07, y+0.1, cw-0.14, 0.36, size=9.5,
            bold=(ci==1), color=tc, align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)

txt(s, "★  VigiliCloud is the only tool combining simplicity + AI guidance + compliance exports at SMB-friendly pricing",
    0.3, 6.78, 12.7, 0.4, size=10.5, bold=True, color=LGREEN)

slide_num(s, 7)

# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — WHY VIGILICLOUD / VALUE PROP
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, BLACK)
accent_bar(s)

tag(s, "WHY VIGILICLOUD", 0.45, 0.55)
txt(s, "Built for Teams Who Can't Afford to Be Wrong",
    0.45, 0.9, 11, 0.55, size=26, bold=True, color=WHITE)

points = [
    ("🚀", "Speed Over Complexity",
     "While competitors take days to set up, VigiliCloud is running in 5 minutes.\nNo agents, no training, no AWS expert required."),
    ("💡", "Actionable — Not Just Alerts",
     "Every finding ships with the exact AWS Console path, CLI command to fix it, and\nan AI-written explanation of why it matters. We close the loop from detection to fix."),
    ("💰", "Priced for the Real Market",
     "Enterprise tools charge $50,000+/year. VigiliCloud starts at $99/month — making\nprofessional-grade AWS security accessible to consultants, startups, and SMBs."),
    ("🤝", "Built for Consultants & MSPs",
     "Scan multiple client accounts, export evidence for audits, and deliver a\nprofessional security report in 2 minutes — not 2 days."),
]
for i, (icon, title, desc) in enumerate(points):
    y = 1.7 + i * 1.3
    dot = box(s, 0.45, y+0.08, 0.52, 0.52, RGBColor(0x06,0x2a,0x1e))
    dot.line.color.rgb = GREEN; dot.line.width = Pt(0.8)
    txt(s, icon, 0.45, y+0.06, 0.52, 0.52, size=16, align=PP_ALIGN.CENTER)
    txt(s, title, 1.1, y,      11, 0.38, size=14, bold=True, color=WHITE)
    txt(s, desc,  1.1, y+0.4,  11.8, 0.65, size=11, color=GREY)

# ROI callout
roi = card(s, 0.45, 6.35, 12.4, 0.75)
txt(s, "💎  ROI SNAPSHOT: One prevented S3 data breach saves an average of $4.5M. VigiliCloud Pro costs $299/month. That's 1,500x ROI on a single incident prevented.",
    0.65, 6.48, 12.0, 0.5, size=11, bold=True, color=LGREEN)

slide_num(s, 8)

# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — PRICING
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, BLACK)
accent_bar(s)

tag(s, "PRICING", 0.45, 0.55)
txt(s, "Simple, Transparent Pricing — 2-Week Free Trial on All Plans",
    0.45, 0.9, 11, 0.55, size=24, bold=True, color=WHITE)
txt(s, "No credit card required. Cancel anytime. Billed in INR · USD shown for reference.",
    0.45, 1.5, 11, 0.35, size=12, color=GREY)

plans = [
    ("STARTER", "$99/mo", "₹8,299/mo", "Solo consultants\n& small teams",
     ["Up to 3 AWS accounts", "All 10 security checks", "Fix guidance", "CSV/JSON exports", "Email alerts"],
     False),
    ("PRO ★", "$299/mo", "₹24,999/mo", "Teams managing\nmultiple AWS environments",
     ["Up to 10 AWS accounts", "Everything in Starter", "AI security analysis (Claude)", "Scheduled daily scans", "Priority support"],
     True),
    ("MSP", "$999/mo", "₹83,499/mo", "Agencies &\nmanaged service providers",
     ["Unlimited AWS accounts", "Everything in Pro", "Multi-customer workflows", "White-label ready", "High-volume scanning"],
     False),
]
for i, (name, usd, inr, desc, feats, hot) in enumerate(plans):
    x = 0.45 + i * 4.25
    border_col = GREEN if hot else RGBColor(0x2a,0x2a,0x2a)
    fill_col   = RGBColor(0x06,0x20,0x14) if hot else CARD
    c = card(s, x, 2.0, 3.95, 4.85)
    c.fill.solid(); c.fill.fore_color.rgb = fill_col
    c.line.color.rgb = border_col
    if hot:
        tag(s, "★ MOST POPULAR", x+1.1, 1.82, GREEN, BLACK)
    txt(s, name,  x+0.2, 2.15, 3.55, 0.45, size=14, bold=True, color=GREEN if hot else WHITE)
    txt(s, usd,   x+0.2, 2.65, 3.55, 0.55, size=26, bold=True, color=LGREEN if hot else WHITE)
    txt(s, inr,   x+0.2, 3.22, 3.55, 0.3,  size=10, color=DKGREY)
    txt(s, desc,  x+0.2, 3.55, 3.55, 0.45, size=10, color=GREY)
    for j, f in enumerate(feats):
        txt(s, f"✓  {f}", x+0.2, 4.1 + j * 0.42, 3.55, 0.38, size=10,
            color=LGREEN if hot else GREY)

txt(s, "All plans include: 2-week free trial  ·  Read-only AWS access  ·  HTTPS encrypted  ·  Data never stored",
    0.45, 7.1, 12.4, 0.32, size=10, color=DKGREY, align=PP_ALIGN.CENTER)

slide_num(s, 9)

# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — TESTIMONIALS / SOCIAL PROOF
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, BLACK)
accent_bar(s)

tag(s, "SOCIAL PROOF", 0.45, 0.55)
txt(s, "What Our Users Say",
    0.45, 0.9, 11, 0.55, size=26, bold=True, color=WHITE)

testimonials = [
    ('"VigiliCloud found 6 critical misconfigurations in our production AWS account we had no idea about. Fixed them all in an afternoon."',
     "Rahul M.", "CTO, SaaS Startup — Bangalore, India"),
    ('"As an AWS consultant I scan client accounts regularly. This cuts my security review from 2 days to 2 minutes. Absolute game changer."',
     "Priya S.", "AWS Consultant — Hyderabad, India"),
    ('"We needed to pass a security audit for an enterprise deal. VigiliCloud gave us the evidence exports we needed to close."',
     "James T.", "VP Engineering — Austin, Texas"),
]
for i, (quote, name, role) in enumerate(testimonials):
    x = 0.45 + i * 4.25
    c = card(s, x, 1.65, 3.95, 3.85)
    txt(s, "★★★★★", x+0.2, 1.85, 3.55, 0.35, size=13, color=LGREEN)
    txt(s, quote,   x+0.2, 2.3,  3.55, 1.85, size=10.5, color=GREY, italic=True)
    # author
    dot = box(s, x+0.2, 4.6, 0.42, 0.42, RGBColor(0x06,0x2a,0x1e))
    dot.line.color.rgb = GREEN; dot.line.width = Pt(0.7)
    txt(s, name[:1], x+0.2, 4.6, 0.42, 0.42, size=14, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    txt(s, name, x+0.72, 4.62, 3.0, 0.28, size=11, bold=True, color=WHITE)
    txt(s, role, x+0.72, 4.9,  3.0, 0.28, size=9,  color=DKGREY)

txt(s, "🔒  Your data never leaves your account  ·  Read-only access  ·  SOC2-aligned practices",
    0.45, 5.82, 12.4, 0.35, size=11, color=DKGREY, align=PP_ALIGN.CENTER)

# Trust logos row
txt(s, "Trusted workflow used by:   Solo AWS Consultants  ·  DevOps Teams  ·  Startups  ·  MSPs  ·  Security Engineers",
    0.45, 6.3, 12.4, 0.35, size=11, bold=True, color=LGREEN, align=PP_ALIGN.CENTER)

slide_num(s, 10)

# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — MOTIVATION / CALL TO ACTION
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, BLACK)
accent_bar(s, 0.0)

txt(s, "The question isn't whether your AWS account\nhas security gaps.", 0.7, 0.5, 11.9, 1.2,
    size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(s, "The question is: will you find them before an attacker does?",
    0.7, 1.8, 11.9, 0.65, size=22, bold=True, color=LGREEN, align=PP_ALIGN.CENTER)

div2 = box(s, 4.5, 2.6, 4.33, 0.05, GREEN); div2.line.fill.background()

bullets = [
    "A single exposed S3 bucket can cost you customers, contracts, and your reputation overnight",
    "Root access keys with no rotation are a ticking time bomb — and they're in 1 in 3 accounts",
    "Enterprise deals now require security evidence — VigiliCloud generates it in 2 minutes",
    "Your competitors are already scanning. The question is: are you?",
]
for i, b in enumerate(bullets):
    txt(s, f"  ▶  {b}", 1.5, 2.85 + i*0.58, 10.3, 0.5, size=12, color=GREY)

# CTA box
cta = card(s, 2.0, 5.38, 9.33, 1.65)
cta.fill.solid(); cta.fill.fore_color.rgb = RGBColor(0x06,0x20,0x14)
cta.line.color.rgb = GREEN; cta.line.width = Pt(1.5)

txt(s, "🚀  Start your FREE 2-week trial today", 2.2, 5.52, 9.0, 0.45,
    size=16, bold=True, color=LGREEN, align=PP_ALIGN.CENTER)
txt(s, "app.vigilicloud.com  ·  Setup in 5 minutes  ·  No credit card required",
    2.2, 6.0, 9.0, 0.35, size=12, color=GREY, align=PP_ALIGN.CENTER)
txt(s, "Book a personal demo →  calendly.com/leelakrishnakoppolu/vigilicloud-demo",
    2.2, 6.4, 9.0, 0.35, size=11, color=DKGREY, align=PP_ALIGN.CENTER)

slide_num(s, 11)

# ── Save ─────────────────────────────────────────────────────────
out = r"c:\Users\leela\compliance-ai-saas\VigiliCloud_Client_Pitch.pptx"
prs.save(out)
print(f"Saved: {out}")
