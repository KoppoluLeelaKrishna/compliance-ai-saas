import sys
from pptx import Presentation

prs = Presentation(r'c:\Users\leela\compliance-ai-saas\AWS_Compliance_AI_SaaS_AWSonly_FullPlan.pptx')

for i, slide in enumerate(prs.slides, 1):
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(t)
    print(f"\n{'='*60}")
    print(f"SLIDE {i}")
    print('='*60)
    for t in texts:
        safe = t.encode('ascii', 'replace').decode()
        print(f"  {safe}")
